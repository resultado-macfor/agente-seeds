import os
from anthropic import Anthropic
import streamlit as st
import io
import google.generativeai as genai
from PIL import Image
import datetime
from openai import OpenAI
from pymongo import MongoClient
from bson import ObjectId
import json
import hashlib
from google.genai import types
import PyPDF2
from pptx import Presentation
import docx
import openai
from typing import List, Dict, Tuple
import hashlib
import pandas as pd
import re
from pypdf import PdfReader, PdfWriter
from pypdf.annotations import Text
import requests
from google.genai import types
import PyPDF2
from pptx import Presentation
import docx
import openai
from typing import List, Dict, Tuple
import hashlib
import pandas as pd
import re
from pypdf import PdfReader, PdfWriter
from pypdf.annotations import Text
import requests
import pdfplumber
from pathlib import Path
import uuid

# Configuração inicial
st.set_page_config(
    layout="wide",
    page_title="Agente Seeds",
    page_icon="🤖"
)

# --- CARREGAR TODAS AS VARIÁVEIS DO AMBIENTE ---
# API Keys
perp_api_key = os.getenv("PERP_API_KEY")
openai_api_key = os.getenv("OPENAI_API_KEY")
anthropic_api_key = os.getenv("ANTHROPIC_API_KEY")
gemini_api_key = os.getenv("GEM_API_KEY")
mistral_api_key = os.getenv("MISTRAL_API_KEY")
rapid_api_key = os.getenv("RAPID_API")
t_api_key = os.getenv("T_API_KEY")
gemini_api_key2 = os.getenv("GEM_API_KEY2")

# Senhas dos squads
senha_admin = os.getenv("SENHA_ADMIN")
senha_syn = os.getenv("SENHA_SYN")
senha_sme = os.getenv("SENHA_SME")
senha_ent = os.getenv("SENHA_ENT")

# MongoDB
mongo_uri = os.getenv("MONGO_URI")

# Astra DB
astra_db_token = os.getenv("ASTRA_DB_APPLICATION_TOKEN")
astra_db_namespace = os.getenv("ASTRA_DB_NAMESPACE")
astra_db_endpoint = os.getenv("ASTRA_DB_API_ENDPOINT")

# --- CONFIGURAÇÃO DOS MODELOS ---
# Configuração da API do Anthropic (Claude)
if anthropic_api_key:
    anthropic_client = Anthropic(api_key=anthropic_api_key)
else:
    st.error("ANTHROPIC_API_KEY não encontrada nas variáveis de ambiente")
    anthropic_client = None

# Configuração da API do Gemini
if gemini_api_key:
    genai.configure(api_key=gemini_api_key)
    modelo_vision = genai.GenerativeModel("gemini-2.5-flash", generation_config={"temperature": 0.0})
    modelo_texto = genai.GenerativeModel("gemini-2.5-flash")
else:
    st.error("GEM_API_KEY não encontrada nas variáveis de ambiente")
    modelo_vision = None
    modelo_texto = None

if openai_api_key:
    openai_client = OpenAI(api_key=openai_api_key)
else:
    st.warning("OPENAI_API_KEY não encontrada nas variáveis de ambiente")
    openai_client = None

if not perp_api_key:
    st.error("PERP_API_KEY não encontrada nas variáveis de ambiente")

# --- Conexão MongoDB ---
if mongo_uri:
    client = MongoClient(mongo_uri)
    db = client['agentes_personalizados']
    collection_agentes = db['agentes']
    collection_conversas = db['conversas']
    collection_usuarios = db['usuarios']
else:
    st.error("MONGO_URI não encontrada nas variáveis de ambiente")
    st.stop()

# --- Sistema de Autenticação ---
def make_hashes(password):
    return hashlib.sha256(str.encode(password)).hexdigest()

def check_hashes(password, hashed_text):
    return make_hashes(password) == hashed_text

# Dados de usuários por squad usando as senhas do ambiente
users_db = {
    "admin@seeds.com": {
        "password": make_hashes(senha_admin if senha_admin else "admin123"),
        "squad": "admin",
        "nome": "Administrador"
    },
    "syngenta@seeds.com": {
        "password": make_hashes(senha_syn if senha_syn else "syn123"),
        "squad": "Syngenta",
        "nome": "Usuário Syngenta"
    },
    "sme@seeds.com": {
        "password": make_hashes(senha_sme if senha_sme else "sme123"),
        "squad": "SME",
        "nome": "Usuário SME"
    },
    "enterprise@seeds.com": {
        "password": make_hashes(senha_ent if senha_ent else "ent123"),
        "squad": "Enterprise",
        "nome": "Usuário Enterprise"
    }
}

def criar_usuario(email, senha, nome, squad):
    """Cria um novo usuário no banco de dados"""
    try:
        if collection_usuarios.find_one({"email": email}):
            return False, "Usuário já existe"
        
        senha_hash = make_hashes(senha)
        
        novo_usuario = {
            "email": email,
            "senha": senha_hash,
            "nome": nome,
            "squad": squad,
            "data_criacao": datetime.datetime.now(),
            "ultimo_login": None,
            "ativo": True
        }
        
        result = collection_usuarios.insert_one(novo_usuario)
        return True, "Usuário criado com sucesso"
        
    except Exception as e:
        return False, f"Erro ao criar usuário: {str(e)}"

def verificar_login(email, senha):
    """Verifica as credenciais do usuário"""
    try:
        usuario = collection_usuarios.find_one({"email": email, "ativo": True})
        
        if usuario:
            if check_hashes(senha, usuario["senha"]):
                collection_usuarios.update_one(
                    {"_id": usuario["_id"]},
                    {"$set": {"ultimo_login": datetime.datetime.now()}}
                )
                return True, usuario, "Login bem-sucedido"
            else:
                return False, None, "Senha incorreta"
        
        if email in users_db:
            user_data = users_db[email]
            if check_hashes(senha, user_data["password"]):
                usuario_fallback = {
                    "email": email,
                    "nome": user_data["nome"],
                    "squad": user_data["squad"],
                    "_id": email
                }
                return True, usuario_fallback, "Login bem-sucedido"
            else:
                return False, None, "Senha incorreta"
        
        return False, None, "Usuário não encontrado"
        
    except Exception as e:
        return False, None, f"Erro no login: {str(e)}"

def get_current_user():
    return st.session_state.get('user', {})

def get_current_squad():
    user = get_current_user()
    return user.get('squad', 'unknown')

def login():
    st.title("🔒 Agente Seeds - Login")
    
    with st.expander("ℹ️ Credenciais de Acesso"):
        st.markdown("""
        **Acesse com as seguintes credenciais:**
        
        | Squad | Email | Senha |
        |-------|-------|-------|
        | Admin | admin@seeds.com | (configurada no .env) |
        | Syngenta | syngenta@seeds.com | (configurada no .env) |
        | SME | sme@seeds.com | (configurada no .env) |
        | Enterprise | enterprise@seeds.com | (configurada no .env) |
        
        *Para criar um novo usuário, utilize a aba de cadastro.*
        """)
    
    tab_login, tab_cadastro = st.tabs(["Login", "Cadastro"])
    
    with tab_login:
        with st.form("login_form"):
            email = st.text_input("Email")
            password = st.text_input("Senha", type="password")
            submit_button = st.form_submit_button("Login")
            
            if submit_button:
                if email and password:
                    sucesso, usuario, mensagem = verificar_login(email, password)
                    if sucesso:
                        st.session_state.logged_in = True
                        st.session_state.user = usuario
                        st.success("Login realizado com sucesso!")
                        st.rerun()
                    else:
                        st.error(mensagem)
                else:
                    st.error("Por favor, preencha todos os campos")
    
    with tab_cadastro:
        with st.form("cadastro_form"):
            st.subheader("Criar Nova Conta")
            
            nome = st.text_input("Nome Completo")
            email = st.text_input("Email")
            squad = st.selectbox(
                "Selecione seu Squad:",
                ["Syngenta", "SME", "Enterprise"],
                help="Escolha o squad ao qual você pertence"
            )
            senha = st.text_input("Senha", type="password")
            confirmar_senha = st.text_input("Confirmar Senha", type="password")
            
            submit_cadastro = st.form_submit_button("Criar Conta")
            
            if submit_cadastro:
                if not all([nome, email, squad, senha, confirmar_senha]):
                    st.error("Por favor, preencha todos os campos")
                elif senha != confirmar_senha:
                    st.error("As senhas não coincidem")
                elif len(senha) < 6:
                    st.error("A senha deve ter pelo menos 6 caracteres")
                else:
                    sucesso, mensagem = criar_usuario(email, senha, nome, squad)
                    if sucesso:
                        st.success("Conta criada com sucesso! Faça login para continuar.")
                    else:
                        st.error(mensagem)

# Verificar login
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False

if not st.session_state.logged_in:
    login()
    st.stop()

# --- RECONFIGURAR APÓS LOGIN ---
if gemini_api_key:
    genai.configure(api_key=gemini_api_key)
    modelo_vision = genai.GenerativeModel("gemini-2.5-flash", generation_config={"temperature": 0.0})
    modelo_texto = genai.GenerativeModel("gemini-2.5-flash")

def check_admin_password():
    return st.session_state.user.get('squad') == "admin"

# --- FUNÇÕES CRUD PARA AGENTES ---
def criar_agente(nome, system_prompt, base_conhecimento, comments, planejamento, categoria, squad_permitido, agente_mae_id=None, herdar_elementos=None):
    agente = {
        "nome": nome,
        "system_prompt": system_prompt,
        "base_conhecimento": base_conhecimento,
        "comments": comments,
        "planejamento": planejamento,
        "categoria": categoria,
        "squad_permitido": squad_permitido,
        "agente_mae_id": agente_mae_id,
        "herdar_elementos": herdar_elementos or [],
        "data_criacao": datetime.datetime.now(),
        "ativo": True,
        "criado_por": get_current_user().get('email', 'unknown'),
        "criado_por_squad": get_current_squad()
    }
    result = collection_agentes.insert_one(agente)
    return result.inserted_id

def listar_agentes():
    current_squad = get_current_squad()
    
    if current_squad == "admin":
        return list(collection_agentes.find({"ativo": True}).sort("data_criacao", -1))
    
    return list(collection_agentes.find({
        "ativo": True,
        "$or": [
            {"squad_permitido": current_squad},
            {"squad_permitido": "Todos"},
            {"criado_por_squad": current_squad}
        ]
    }).sort("data_criacao", -1))

def listar_agentes_para_heranca(agente_atual_id=None):
    current_squad = get_current_squad()
    query = {"ativo": True}
    
    if current_squad != "admin":
        query["$or"] = [
            {"squad_permitido": current_squad},
            {"squad_permitido": "Todos"},
            {"criado_por_squad": current_squad}
        ]
    
    if agente_atual_id:
        if isinstance(agente_atual_id, str):
            agente_atual_id = ObjectId(agente_atual_id)
        query["_id"] = {"$ne": agente_atual_id}
    
    return list(collection_agentes.find(query).sort("data_criacao", -1))

def obter_agente(agente_id):
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    
    agente = collection_agentes.find_one({"_id": agente_id})
    
    if agente and agente.get('ativo', True):
        current_squad = get_current_squad()
        
        if current_squad == "admin":
            return agente
        
        squad_permitido = agente.get('squad_permitido')
        criado_por_squad = agente.get('criado_por_squad')
        
        if squad_permitido == current_squad or squad_permitido == "Todos" or criado_por_squad == current_squad:
            return agente
    
    return None

def atualizar_agente(agente_id, nome, system_prompt, base_conhecimento, comments, planejamento, categoria, squad_permitido, agente_mae_id=None, herdar_elementos=None):
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    
    agente_existente = obter_agente(agente_id)
    if not agente_existente:
        raise PermissionError("Agente não encontrado ou sem permissão de edição")
    
    return collection_agentes.update_one(
        {"_id": agente_id},
        {
            "$set": {
                "nome": nome,
                "system_prompt": system_prompt,
                "base_conhecimento": base_conhecimento,
                "comments": comments,
                "planejamento": planejamento,
                "categoria": categoria,
                "squad_permitido": squad_permitido,
                "agente_mae_id": agente_mae_id,
                "herdar_elementos": herdar_elementos or [],
                "data_atualizacao": datetime.datetime.now()
            }
        }
    )

def desativar_agente(agente_id):
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    
    agente_existente = obter_agente(agente_id)
    if not agente_existente:
        raise PermissionError("Agente não encontrado ou sem permissão para desativar")
    
    return collection_agentes.update_one(
        {"_id": agente_id},
        {"$set": {"ativo": False, "data_desativacao": datetime.datetime.now()}}
    )

def obter_agente_com_heranca(agente_id):
    agente = obter_agente(agente_id)
    if not agente or not agente.get('agente_mae_id'):
        return agente
    
    agente_mae = obter_agente(agente['agente_mae_id'])
    if not agente_mae:
        return agente
    
    elementos_herdar = agente.get('herdar_elementos', [])
    agente_completo = agente.copy()
    
    for elemento in elementos_herdar:
        if elemento == 'system_prompt' and not agente_completo.get('system_prompt'):
            agente_completo['system_prompt'] = agente_mae.get('system_prompt', '')
        elif elemento == 'base_conhecimento' and not agente_completo.get('base_conhecimento'):
            agente_completo['base_conhecimento'] = agente_mae.get('base_conhecimento', '')
        elif elemento == 'comments' and not agente_completo.get('comments'):
            agente_completo['comments'] = agente_mae.get('comments', '')
        elif elemento == 'planejamento' and not agente_completo.get('planejamento'):
            agente_completo['planejamento'] = agente_mae.get('planejamento', '')
    
    return agente_completo

def salvar_conversa(agente_id, mensagens, segmentos_utilizados=None):
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    conversa = {
        "agente_id": agente_id,
        "mensagens": mensagens,
        "segmentos_utilizados": segmentos_utilizados,
        "data_criacao": datetime.datetime.now()
    }
    return collection_conversas.insert_one(conversa)

def obter_conversas(agente_id, limite=10):
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    return list(collection_conversas.find(
        {"agente_id": agente_id}
    ).sort("data_criacao", -1).limit(limite))

def construir_contexto(agente, segmentos_selecionados, historico_mensagens=None):
    contexto = ""
    
    if "system_prompt" in segmentos_selecionados and agente.get('system_prompt'):
        contexto += f"### INSTRUÇÕES DO SISTEMA ###\n{agente['system_prompt']}\n\n"
    
    if "base_conhecimento" in segmentos_selecionados and agente.get('base_conhecimento'):
        contexto += f"### BASE DE CONHECIMENTO ###\n{agente['base_conhecimento']}\n\n"
    
    if "comments" in segmentos_selecionados and agente.get('comments'):
        contexto += f"### Diário DO CLIENTE ###\n{agente['comments']}\n\n"
    
    if "planejamento" in segmentos_selecionados and agente.get('planejamento'):
        contexto += f"### PLANEJAMENTO ###\n{agente['planejamento']}\n\n"
    
    if historico_mensagens:
        contexto += "### HISTÓRICO DA CONVERSA ###\n"
        for msg in historico_mensagens:
            contexto += f"{msg['role']}: {msg['content']}\n"
        contexto += "\n"
    
    contexto += "### RESPOSTA ATUAL ###\nassistant:"
    
    return contexto

def gerar_resposta_modelo(prompt: str, modelo_escolhido: str = "Gemini", contexto_agente: str = None) -> str:
    try:
        if modelo_escolhido == "Gemini" and modelo_texto:
            if contexto_agente:
                prompt_completo = f"{contexto_agente}\n\n{prompt}"
            else:
                prompt_completo = prompt
            
            resposta = modelo_texto.generate_content(prompt_completo)
            return resposta.text
            
        elif modelo_escolhido == "Claude" and anthropic_client:
            if contexto_agente:
                system_prompt = contexto_agente
            else:
                system_prompt = "Você é um assistente útil."
            
            message = anthropic_client.messages.create(
                max_tokens=4000,
                messages=[{"role": "user", "content": prompt}],
                model="claude-haiku-4-5-20251001",
                system=system_prompt
            )
            return message.content[0].text
            
        elif modelo_escolhido == "OpenAI" and openai_client:
            try:
                response = openai_client.responses.create(
                    model="gpt-4o-mini",
                    input=prompt,
                    instructions=contexto_agente if contexto_agente else "Você é um assistente útil."
                )
                return response.output_text
            except Exception:
                messages = []
                if contexto_agente:
                    messages.append({"role": "system", "content": contexto_agente})
                messages.append({"role": "user", "content": prompt})
                
                response = openai_client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=messages,
                    max_tokens=4000,
                    temperature=0.0
                )
                return response.choices[0].message.content
            
        else:
            return f"❌ Modelo {modelo_escolhido} não disponível."
            
    except Exception as e:
        return f"❌ Erro ao gerar resposta com {modelo_escolhido}: {str(e)}"

# --- FUNÇÕES DE BUSCA WEB ---
def realizar_busca_web_com_fontes(termos_busca: str, contexto_agente: str = None) -> str:
    if not perp_api_key:
        return "❌ API do Perplexity não configurada."
    
    try:
        headers = {
            "Authorization": f"Bearer {perp_api_key}",
            "Content-Type": "application/json"
        }
        
        mensagem_sistema = contexto_agente if contexto_agente else "Você é um assistente de pesquisa que fornece informações precisas e atualizadas COM FONTES."
        
        data = {
            "model": "sonar",
            "messages": [
                {
                    "role": "system",
                    "content": f"{mensagem_sistema}\n\nIMPORTANTE: Você DEVE SEMPRE incluir as fontes (links e nomes dos sites) de onde tirou as informações."
                },
                {
                    "role": "user", 
                    "content": f"""Pesquise informações sobre: {termos_busca}

                    REQUISITOS OBRIGATÓRIOS:
                    1. Forneça informações TÉCNICAS e ATUALIZADAS
                    2. INCLUA SEMPRE as fontes para cada informação
                    3. Use o formato: **Fonte: [Nome do Site/Portal] ([link completo])**
                    4. Priorize fontes confiáveis
                    5. Forneça dados concretos: números, estatísticas, resultados"""
                }
            ],
            "max_tokens": 4000,
            "temperature": 0.0
        }
        
        response = requests.post(
            "https://api.perplexity.ai/chat/completions",
            headers=headers,
            json=data,
            timeout=60
        )
        
        if response.status_code == 200:
            result = response.json()
            return result['choices'][0]['message']['content']
        else:
            return f"❌ Erro na busca web: {response.status_code}"
                
    except Exception as e:
        return f"❌ Erro ao realizar busca web: {str(e)}"

def analisar_urls_com_fontes(urls: List[str], pergunta: str, contexto_agente: str = None) -> str:
    try:
        headers = {
            "Authorization": f"Bearer {perp_api_key}",
            "Content-Type": "application/json"
        }
        
        urls_contexto = "\n".join([f"- {url}" for url in urls])
        
        messages = []
        
        if contexto_agente:
            messages.append({
                "role": "system",
                "content": f"Contexto do agente: {contexto_agente}\n\nIMPORTANTE: Sempre cite as fontes específicas das URLs analisadas."
            })
        else:
            messages.append({
                "role": "system",
                "content": "Você é um analista de conteúdo. Sempre cite as fontes específicas das URLs analisadas."
            })
        
        messages.append({
            "role": "user",
            "content": f"""Analise as seguintes URLs e responda à pergunta:

URLs para análise:
{urls_contexto}

Pergunta: {pergunta}

REQUISITOS OBRIGATÓRIOS:
1. Para cada informação, mencione de qual URL específica veio
2. Use formato: **Fonte: [Nome do Site/Portal] ([URL específica])**"""
        })
        
        data = {
            "model": "sonar-medium-online",
            "messages": messages,
            "max_tokens": 3000,
            "temperature": 0.0
        }
        
        response = requests.post(
            "https://api.perplexity.ai/chat/completions",
            headers=headers,
            json=data,
            timeout=45
        )
        
        if response.status_code == 200:
            result = response.json()
            return result['choices'][0]['message']['content']
        else:
            return f"❌ Erro na análise: {response.status_code}"
                
    except Exception as e:
        return f"❌ Erro ao analisar URLs: {str(e)}"

# --- FUNÇÕES DE EXTRAÇÃO DE TEXTO ---
def extract_text_from_pdf_com_slides(arquivo_pdf):
    try:
        pdf_reader = PyPDF2.PdfReader(arquivo_pdf)
        slides_info = []
        
        for pagina_num, pagina in enumerate(pdf_reader.pages):
            texto = pagina.extract_text()
            slides_info.append({
                'numero': pagina_num + 1,
                'conteudo': texto,
                'tipo': 'página'
            })
        
        texto_completo = "\n\n".join([f"--- PÁGINA {s['numero']} ---\n{s['conteudo']}" for s in slides_info])
        return texto_completo, slides_info
        
    except Exception as e:
        return f"Erro na extração PDF: {str(e)}", []

def extract_text_from_pptx_com_slides(arquivo_pptx):
    try:
        prs = Presentation(io.BytesIO(arquivo_pptx.read()))
        slides_info = []
        
        for slide_num, slide in enumerate(prs.slides):
            texto_slide = f"--- SLIDE {slide_num + 1} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texto_slide += shape.text + "\n"
            
            slides_info.append({
                'numero': slide_num + 1,
                'conteudo': texto_slide,
                'tipo': 'slide'
            })
        
        texto_completo = "\n\n".join([s['conteudo'] for s in slides_info])
        return texto_completo, slides_info
        
    except Exception as e:
        return f"Erro na extração PPTX: {str(e)}", []

def extrair_texto_arquivo(arquivo):
    try:
        if arquivo.type == "text/plain":
            return str(arquivo.read(), "utf-8")
        elif arquivo.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = docx.Document(io.BytesIO(arquivo.read()))
            texto = ""
            for para in doc.paragraphs:
                texto += para.text + "\n"
            return texto
        else:
            return f"Tipo não suportado: {arquivo.type}"
    except Exception as e:
        return f"Erro na extração: {str(e)}"

def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
    except Exception as e:
        pass

    if len(text.strip()) < 100:
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text 
        except Exception as e:
            pass

    return text

# --- FUNÇÕES DE ANÁLISE DE TEXTO ---
def extrair_score(texto_analise):
    padrao = r'SCORE.*?\[(\d+)(?:/10)?\]'
    correspondencias = re.findall(padrao, texto_analise, re.IGNORECASE)
    if correspondencias:
        return int(correspondencias[0])
    return 5

def criar_analisadores_texto(contexto_agente, contexto_global):
    analisadores = {
        'ortografia': {
            'nome': '🔤 Especialista em Ortografia e Gramática',
            'prompt': f"""
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM ORTOGRAFIA E GRAMÁTICA PORTUGUÊS BR

**Sua tarefa:** Analisar EXCLUSIVAMENTE aspectos ortográficos e gramaticais.

### CRITÉRIOS DE ANÁLISE:
1. **Ortografia** - Erros de escrita
2. **Gramática** - Concordância, regência, colocação
3. **Pontuação** - Uso de vírgulas, pontos, etc.
4. **Acentuação** - Erros de acentuação

### FORMATO DE RESPOSTA OBRIGATÓRIO:

## 🔤 RELATÓRIO ORTOGRÁFICO

### ✅ ACERTOS
- [Itens corretos]

### ❌ ERROS IDENTIFICADOS
- [Lista específica de erros com correções]

### 📊 SCORE ORTOGRÁFICO: [X/10]
"""
        },
        'lexico': {
            'nome': '📚 Especialista em Léxico e Vocabulário',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM LÉXICO E VOCABULÁRIO

### CRITÉRIOS DE ANÁLISE:
1. **Variedade Lexical** - Riqueza de vocabulário
2. **Precisão Semântica** - Uso adequado das palavras
3. **Repetição** - Palavras ou expressões repetidas em excesso

### FORMATO DE RESPOSTA:

## 📚 RELATÓRIO LEXICAL

### ✅ VOCABULÁRIO ADEQUADO
- [Pontos fortes do vocabulário]

### ⚠️ ASPECTOS A MELHORAR
- [Problemas lexicais identificados]

### 📊 SCORE LEXICAL: [X/10]
"""
        },
        'branding': {
            'nome': '🎨 Especialista em Branding e Identidade',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM BRANDING E IDENTIDADE

### CRITÉRIOS DE ANÁLISE:
1. **Tom de Voz** - Alinhamento com personalidade da marca
2. **Mensagem Central** - Consistência da mensagem
3. **Valores da Marca** - Reflexo dos valores organizacionais

### FORMATO DE RESPOSTA:

## 🎨 RELATÓRIO DE BRANDING

### ✅ ALINHAMENTOS
- [Elementos que seguem as diretrizes]

### ❌ DESVIOS IDENTIFICADOS
- [Elementos fora do padrão da marca]

### 📊 SCORE BRANDING: [X/10]
"""
        }
    }
    return analisadores

def executar_analise_texto_especializada(texto, nome_arquivo, analisadores):
    resultados = {}
    
    for area, config in analisadores.items():
        with st.spinner(f"Executando {config['nome']}..."):
            try:
                prompt_completo = f"""
{config['prompt']}

###BEGIN TEXTO PARA ANÁLISE###
**Arquivo:** {nome_arquivo}
**Conteúdo:**
{texto[:8000]}
###END TEXTO PARA ANÁLISE###

Por favor, forneça sua análise no formato solicitado.
"""
                
                resposta = modelo_texto.generate_content(prompt_completo)
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': resposta.text,
                    'score': extrair_score(resposta.text)
                }
                
            except Exception as e:
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': f"❌ Erro na análise: {str(e)}",
                    'score': 0
                }
    
    return resultados

def gerar_relatorio_texto_consolidado(resultados_especialistas, nome_arquivo):
    relatorio = f"""
# 📊 RELATÓRIO CONSOLIDADO DE VALIDAÇÃO

**Documento:** {nome_arquivo}
**Data da Análise:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}

## 🎖️ SCORES POR ÁREA
"""
    
    for area, resultado in resultados_especialistas.items():
        emoji = "✅" if resultado['score'] >= 8 else "⚠️" if resultado['score'] >= 6 else "❌"
        relatorio += f"- {emoji} **{resultado['nome']}:** {resultado['score']}/10\n"
    
    relatorio += "\n## 📋 ANÁLISES DETALHADAS POR ESPECIALISTA\n"
    
    for area, resultado in resultados_especialistas.items():
        relatorio += f"\n### {resultado['nome']}\n"
        relatorio += f"{resultado['analise']}\n"
        relatorio += "---\n"
    
    return relatorio

# --- FUNÇÕES DE REVISÃO ORTOGRÁFICA ---
def revisar_texto_ortografia(texto, agente, segmentos_selecionados, revisao_estilo=True, manter_estrutura=True, explicar_alteracoes=True, modelo_escolhido="Gemini"):
    contexto_agente = "CONTEXTO DO AGENTE PARA REVISÃO:\n\n"
    
    if "system_prompt" in segmentos_selecionados and agente.get('system_prompt'):
        contexto_agente += f"DIRETRIZES PRINCIPAIS:\n{agente['system_prompt']}\n\n"
    
    if "base_conhecimento" in segmentos_selecionados and agente.get('base_conhecimento'):
        contexto_agente += f"BASE DE CONHECIMENTO:\n{agente['base_conhecimento']}\n\n"
    
    if "comments" in segmentos_selecionados and agente.get('comments'):
        contexto_agente += f"COMENTÁRIOS E OBSERVAÇÕES:\n{agente['comments']}\n\n"
    
    instrucoes_revisao = ""
    
    if revisao_estilo:
        instrucoes_revisao += """
        - Analise e melhore a clareza, coesão e coerência textual
        - Verifique adequação ao tom da marca
        - Elimine vícios de linguagem e redundâncias
        """
    
    if manter_estrutura:
        instrucoes_revisao += """
        - Mantenha a estrutura geral do texto original
        - Preserve parágrafos e seções quando possível
        """
    
    if explicar_alteracoes:
        instrucoes_revisao += """
        - Inclua justificativa para as principais alterações
        - Explique correções gramaticais importantes
        """
    
    prompt_revisao = f"""
{contexto_agente}

TEXTO PARA REVISÃO:
{texto}

INSTRUÇÕES PARA REVISÃO:

1. **REVISÃO ORTOGRÁFICA E GRAMATICAL:**
   - Corrija erros de ortografia, acentuação e grafia
   - Verifique concordância nominal e verbal
   - Ajuste pontuação

2. **REVISÃO DE ESTILO E CLAREZA:**
   {instrucoes_revisao}

FORMATO DA RESPOSTA:

## 📋 TEXTO REVISADO
[Texto completo revisado]

## 🔍 PRINCIPAIS ALTERAÇÕES REALIZADAS
[Lista das principais correções]

## 📊 RESUMO DA REVISÃO
[Resumo dos problemas encontrados]
"""
    
    try:
        resposta = gerar_resposta_modelo(prompt_revisao, modelo_escolhido)
        return resposta
    except Exception as e:
        return f"❌ Erro durante a revisão: {str(e)}"

# --- FUNÇÕES PARA VALIDAÇÃO DE TEXTO EM IMAGEM ---
def gerar_relatorio_texto_imagem_consolidado(resultados):
    relatorio = f"""
# 📝 RELATÓRIO DE VALIDAÇÃO DE TEXTO EM IMAGEM

**Data da Análise:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}
**Total de Imagens Analisadas:** {len(resultados)}

## 📋 ANÁLISE INDIVIDUAL POR ARTE
"""
    
    for resultado in resultados:
        relatorio += f"\n{resultado['analise']}\n"
    
    relatorio += "\n\n## 📌 RESUMO FINAL\n"
    relatorio += "Arte\tErros encontrados?\tObservações\n"
    relatorio += "---\t---\t---\n"
    
    for resultado in resultados:
        status_text = {
            "Correto": "❌ Não",
            "Ajustes sugeridos": "⚠️ Sugestões apenas",
            "Com erros": "✅ Sim",
            "Erro": "❌ Erro na análise"
        }.get(resultado['status'], "❓ Desconhecido")
        
        relatorio += f"Arte {resultado['indice']}\t{status_text}\t{resultado['status']}\n"
    
    return relatorio

# --- FUNÇÕES DE ANÁLISE DE IMAGEM ---
def criar_analisadores_imagem(contexto_agente, contexto_global):
    analisadores = {
        'composicao_visual': {
            'nome': '🎨 Especialista em Composição Visual',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM COMPOSIÇÃO VISUAL

### CRITÉRIOS DE ANÁLISE:
1. **Balanceamento** - Distribuição equilibrada dos elementos
2. **Hierarquia Visual** - Foco e pontos de atenção
3. **Harmonia** - Conjunto visual coeso

### FORMATO DE RESPOSTA:

## 🎨 RELATÓRIO DE COMPOSIÇÃO VISUAL

### ✅ PONTOS FORTES
- [Elementos bem compostos]

### ⚠️ PROBLEMAS DE COMPOSIÇÃO
- [Issues de organização visual]

### 📊 SCORE COMPOSIÇÃO: [X/10]
"""
        },
        'cores_branding': {
            'nome': '🌈 Especialista em Cores e Branding',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUNÇÃO: ESPECIALISTA EM CORES E BRANDING

### CRITÉRIOS DE ANÁLISE:
1. **Paleta de Cores** - Cores utilizadas
2. **Contraste** - Legibilidade e visibilidade
3. **Consistência** - Coerência com identidade visual

### FORMATO DE RESPOSTA:

## 🌈 RELATÓRIO DE CORES

### ✅ CORES ALINHADAS
- [Cores que seguem as diretrizes]

### ❌ PROBLEMAS DE COR
- [Cores fora do padrão]

### 📊 SCORE CORES: [X/10]
"""
        }
    }
    return analisadores

def executar_analise_imagem_especializada(uploaded_image, nome_imagem, analisadores):
    resultados = {}
    
    for area, config in analisadores.items():
        with st.spinner(f"Executando {config['nome']}..."):
            try:
                prompt_completo = f"""
{config['prompt']}

###BEGIN IMAGEM PARA ANÁLISE###
**Arquivo:** {nome_imagem}
###END IMAGEM PARA ANÁLISE###

Por favor, forneça sua análise no formato solicitado.
"""
                
                response = modelo_vision.generate_content([
                    prompt_completo,
                    {"mime_type": "image/jpeg", "data": uploaded_image.getvalue()}
                ])
                
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': response.text,
                    'score': extrair_score(response.text)
                }
                
            except Exception as e:
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': f"❌ Erro na análise: {str(e)}",
                    'score': 0
                }
    
    return resultados

def gerar_relatorio_imagem_consolidado(resultados_especialistas, nome_imagem, dimensoes):
    relatorio = f"""
# 🖼️ RELATÓRIO CONSOLIDADO DE IMAGEM

**Arquivo:** {nome_imagem}
**Dimensões:** {dimensoes}
**Data da Análise:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}

## 🎖️ SCORES POR ÁREA ESPECIALIZADA
"""
    
    for area, resultado in resultados_especialistas.items():
        emoji = "✅" if resultado['score'] >= 8 else "⚠️" if resultado['score'] >= 6 else "❌"
        relatorio += f"- {emoji} **{resultado['nome']}:** {resultado['score']}/10\n"
    
    relatorio += "\n## 📋 ANÁLISES DETALHADAS\n"
    
    for area, resultado in resultados_especialistas.items():
        relatorio += f"\n### {resultado['nome']}\n"
        relatorio += f"{resultado['analise']}\n"
        relatorio += "---\n"
    
    return relatorio

# --- FUNÇÕES PARA ANÁLISE DE VÍDEO ---
def criar_analisadores_video(contexto_agente, contexto_global, contexto_video_especifico):
    analisadores = {
        'narrativa_estrutura': {
            'nome': '📖 Especialista em Narrativa e Estrutura',
            'prompt': f"""
{contexto_agente}
{contexto_global}
{contexto_video_especifico}

## FUNÇÃO: ESPECIALISTA EM NARRATIVA E ESTRUTURA

### CRITÉRIOS DE ANÁLISE:
1. **Arco Narrativo** - Desenvolvimento da história
2. **Ritmo** - Velocidade e fluidez da narrativa
3. **Estrutura** - Organização do conteúdo

### FORMATO DE RESPOSTA:

## 📖 RELATÓRIO DE NARRATIVA

### ✅ PONTOS FORTES
- [Elementos narrativos bem executados]

### ⚠️ PROBLEMAS DE ESTRUTURA
- [Issues na organização]

### 📊 SCORE NARRATIVA: [X/10]
"""
        },
        'qualidade_audio': {
            'nome': '🔊 Especialista em Qualidade de Áudio',
            'prompt': f"""
{contexto_agente}
{contexto_global}
{contexto_video_especifico}

## FUNÇÃO: ESPECIALISTA EM QUALIDADE DE ÁUDIO

### CRITÉRIOS DE ANÁLISE:
1. **Clareza Vocal** - Inteligibilidade da fala
2. **Qualidade Técnica** - Ruído, distorção
3. **Sincronização** - Relação áudio-vídeo

### FORMATO DE RESPOSTA:

## 🔊 RELATÓRIO DE ÁUDIO

### ✅ ACERTOS
- [Elementos sonoros bem executados]

### ❌ PROBLEMAS
- [Issues técnicos]

### 📊 SCORE ÁUDIO: [X/10]
"""
        }
    }
    return analisadores

def executar_analise_video_especializada(uploaded_video, nome_video, analisadores):
    resultados = {}
    
    for area, config in analisadores.items():
        with st.spinner(f"Executando {config['nome']}..."):
            try:
                prompt_completo = f"""
{config['prompt']}

###BEGIN VÍDEO PARA ANÁLISE###
**Arquivo:** {nome_video}
###END VÍDEO PARA ANÁLISE###

Por favor, forneça sua análise no formato solicitado.
"""
                
                video_bytes = uploaded_video.getvalue()
                
                response = modelo_vision.generate_content([
                    prompt_completo,
                    {"mime_type": uploaded_video.type, "data": video_bytes}
                ])
                
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': response.text,
                    'score': extrair_score(response.text)
                }
                
            except Exception as e:
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': f"❌ Erro na análise: {str(e)}",
                    'score': 0
                }
    
    return resultados

def gerar_relatorio_video_consolidado(resultados_especialistas, nome_video, tipo_video):
    relatorio = f"""
# 🎬 RELATÓRIO CONSOLIDADO DE VÍDEO

**Arquivo:** {nome_video}
**Formato:** {tipo_video}
**Data da Análise:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}

## 🎖️ SCORES POR ÁREA ESPECIALIZADA
"""
    
    for area, resultado in resultados_especialistas.items():
        emoji = "✅" if resultado['score'] >= 8 else "⚠️" if resultado['score'] >= 6 else "❌"
        relatorio += f"- {emoji} **{resultado['nome']}:** {resultado['score']}/10\n"
    
    relatorio += "\n## 📋 ANÁLISES DETALHADAS\n"
    
    for area, resultado in resultados_especialistas.items():
        relatorio += f"\n### {resultado['nome']}\n"
        relatorio += f"{resultado['analise']}\n"
        relatorio += "---\n"
    
    return relatorio

# --- FUNÇÕES PARA COMENTÁRIOS EM PDF ---
def extrair_comentarios_analise(texto_analise):
    comentarios = []
    
    padroes = [
        r'❌\s*(.*?)(?=\n|$)',
        r'⚠️\s*(.*?)(?=\n|$)',
        r'PROBLEMAS.*?\n(.*?)(?=###|\n\n|$)',
        r'DESVIOS.*?\n(.*?)(?=###|\n\n|$)'
    ]
    
    for padrao in padroes:
        matches = re.findall(padrao, texto_analise, re.IGNORECASE | re.DOTALL)
        for match in matches:
            if isinstance(match, tuple):
                match = match[0]
            comentario = match.strip()
            if comentario and len(comentario) > 10:
                comentarios.append(comentario)
    
    return comentarios[:10]

def adicionar_comentarios_pdf(arquivo_pdf_original, comentarios, nome_documento):
    try:
        reader = PdfReader(io.BytesIO(arquivo_pdf_original.getvalue()))
        writer = PdfWriter()
        
        for page in reader.pages:
            writer.add_page(page)
        
        for i, comentario in enumerate(comentarios):
            if i >= 5:
                break
                
            y_pos = 750 - (i * 100)
            
            annotation = Text(
                text=f"📝 Comentário {i+1}: {comentario[:200]}...",
                rect=(50, y_pos, 400, y_pos + 20),
                open=False
            )
            
            writer.add_annotation(page_number=0, annotation=annotation)
        
        pdf_com_comentarios = io.BytesIO()
        writer.write(pdf_com_comentarios)
        pdf_com_comentarios.seek(0)
        
        return pdf_com_comentarios
        
    except Exception as e:
        st.error(f"❌ Erro ao adicionar comentários ao PDF: {str(e)}")
        return None

def criar_relatorio_comentarios(comentarios, nome_documento, contexto_analise):
    relatorio = f"""
# 📋 RELATÓRIO DE COMENTÁRIOS - {nome_documento}

**Data da Análise:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}
**Total de Comentários:** {len(comentarios)}

## 📝 COMENTÁRIOS E SUGESTÕES
"""
    
    for i, comentario in enumerate(comentarios, 1):
        relatorio += f"### 🔍 Comentário {i}\n{comentario}\n\n"
    
    return relatorio

# --- FUNÇÕES PARA GERAÇÃO DE CONTEÚDO ---
def gerar_conteudo_modelo(prompt: str, modelo_escolhido: str = "Gemini", contexto_agente: str = None) -> str:
    try:
        if modelo_escolhido == "Gemini" and modelo_texto:
            if contexto_agente:
                prompt_completo = f"{contexto_agente}\n\n{prompt}"
            else:
                prompt_completo = prompt
            
            resposta = modelo_texto.generate_content(prompt_completo)
            return resposta.text
            
        elif modelo_escolhido == "Claude" and anthropic_client:
            if contexto_agente:
                system_prompt = contexto_agente
            else:
                system_prompt = "Você é um assistente útil para geração de conteúdo."
            
            message = anthropic_client.messages.create(
                max_tokens=4000,
                messages=[{"role": "user", "content": prompt}],
                model="claude-haiku-4-5-20251001",
                system=system_prompt
            )
            return message.content[0].text
            
        elif modelo_escolhido == "OpenAI" and openai_client:
            try:
                response = openai_client.responses.create(
                    model="gpt-4o-mini",
                    input=prompt,
                    instructions=contexto_agente if contexto_agente else "Você é um assistente especializado em geração de conteúdo."
                )
                return response.output_text
            except Exception:
                messages = []
                if contexto_agente:
                    messages.append({"role": "system", "content": contexto_agente})
                messages.append({"role": "user", "content": prompt})
                
                response = openai_client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=messages,
                    max_tokens=4000,
                    temperature=0.0
                )
                return response.choices[0].message.content
            
        else:
            return f"❌ Modelo {modelo_escolhido} não disponível"
            
    except Exception as e:
        return f"❌ Erro ao gerar conteúdo: {str(e)}"

def transcrever_audio_video(arquivo, tipo):
    return f"Transcrição do {tipo} {arquivo.name} - Esta funcionalidade requer configuração adicional de APIs de transcrição."

# --- FUNÇÕES PARA AGENTE DE MONITORAMENTO ---
def gerar_resposta_agente(pergunta_usuario: str, historico: List[Dict] = None, agente_monitoramento=None, modelo_escolhido="Gemini", contexto_adicional: str = None) -> str:
    if agente_monitoramento and agente_monitoramento.get('base_conhecimento'):
        system_prompt = agente_monitoramento['base_conhecimento']
    else:
        system_prompt = """
        PERSONALIDADE: Especialista com habilidade social - "Especialista que fala como gente"

        TOM DE VOZ:
        - Técnico, confiável e seguro, mas acessível
        - Evita exageros e promessas vazias
        - Sempre embasado em fatos e ciência
        - Frases curtas e diretas, mais simpáticas
        - Toque de leveza e ironia pontual quando o contexto permite
        """

    contexto_completo = system_prompt
    if contexto_adicional and contexto_adicional.strip():
        contexto_completo += f"\n\nCONTEXTO ADICIONAL:\n{contexto_adicional}"
    
    prompt_final = f"""
{contexto_completo}

PERGUNTA DO USUÁRIO:
{pergunta_usuario}

INSTRUÇÕES:
- Adapte seu tom ao tipo de pergunta
- Forneça respostas breves - 1 a 2 frases
- Não invente informações técnicas
- Responda de forma direta e objetiva
"""
    
    try:
        resposta = gerar_resposta_modelo(prompt_final, modelo_escolhido)
        return resposta
    except Exception as e:
        return f"Erro ao gerar resposta: {str(e)}"

# --- INICIALIZAÇÃO DE SESSION STATE ---
if 'analise_especializada_texto' not in st.session_state:
    st.session_state.analise_especializada_texto = True
if 'analise_especializada_imagem' not in st.session_state:
    st.session_state.analise_especializada_imagem = True
if 'analise_especializada_video' not in st.session_state:
    st.session_state.analise_especializada_video = True
if 'analisadores_selecionados_texto' not in st.session_state:
    st.session_state.analisadores_selecionados_texto = ['ortografia', 'lexico', 'branding']
if 'analisadores_selecionados_imagem' not in st.session_state:
    st.session_state.analisadores_selecionados_imagem = ['composicao_visual', 'cores_branding']
if 'analisadores_selecionados_video' not in st.session_state:
    st.session_state.analisadores_selecionados_video = ['narrativa_estrutura', 'qualidade_audio']
if 'validacao_triggered' not in st.session_state:
    st.session_state.validacao_triggered = False
if 'todos_textos' not in st.session_state:
    st.session_state.todos_textos = []
if 'resultados_analise_imagem' not in st.session_state:
    st.session_state.resultados_analise_imagem = []
if 'resultados_analise_video' not in st.session_state:
    st.session_state.resultados_analise_video = []
if 'conteudo_gerado' not in st.session_state:
    st.session_state.conteudo_gerado = None
if 'calendario_gerado' not in st.session_state:
    st.session_state.calendario_gerado = None
if 'agente_selecionado' not in st.session_state:
    st.session_state.agente_selecionado = None
if 'messages' not in st.session_state:
    st.session_state.messages = []
if 'segmentos_selecionados' not in st.session_state:
    st.session_state.segmentos_selecionados = []
if 'messages_monitoramento' not in st.session_state:
    st.session_state.messages_monitoramento = []

# --- DICIONÁRIO DE PRODUTOS PARA BRIEFING ---
PRODUCT_DESCRIPTIONS = {
    "megafol": "Bioativador natural que potencializa o metabolismo da planta",
    "verdatis": "Inseticida com tecnologia PLINAZOLIN para controle de pragas",
    "fortenza": "Tratamento de sementes inseticida para proteção inicial",
    "miravis duo": "Fungicida para controle de manchas foliares",
    "engeo pleno s": "Inseticida para controle de percevejos",
    "axial": "Herbicida para controle de gramíneas",
    "elestal neo": "Herbicida para controle de plantas daninhas",
    "reverb": "Fungicida para controle de doenças",
    "certano hf": "Herbicida para controle de plantas daninhas"
}

def extract_product_info(content):
    """Extrai produto, cultura e ação do conteúdo da célula"""
    content_lower = content.lower()
    
    for product in PRODUCT_DESCRIPTIONS.keys():
        if product in content_lower:
            parts = content_lower.split('-')
            product_found = product
            
            culture = ""
            action = ""
            
            if len(parts) >= 2:
                culture = parts[1].strip()
            if len(parts) >= 3:
                action = parts[2].strip()
            
            return product_found, culture, action
    
    return None, None, None

def generate_briefing(content, product_name, culture, action, data_input, formato_principal):
    """Gera briefing completo no padrão SYN"""
    
    meses = {
        1: "janeiro", 2: "fevereiro", 3: "março", 4: "abril",
        5: "maio", 6: "junho", 7: "julho", 8: "agosto",
        9: "setembro", 10: "outubro", 11: "novembro", 12: "dezembro"
    }
    mes = meses[data_input.month]
    
    descricao_produto = PRODUCT_DESCRIPTIONS.get(product_name, "Produto SYN")
    
    briefing = f"""
# BRIEFING DE CONTEÚDO - SYN

## INFORMAÇÕES GERAIS
**Produto:** {product_name.upper()}
**Cultura:** {culture if culture else "Não especificada"}
**Ação/Tema:** {action if action else content}
**Data de Publicação:** {mes}/{data_input.year}
**Formato Principal:** {formato_principal}

## DESCRIÇÃO DO PRODUTO
{descricao_produto}

## OBJETIVO DO CONTEÚDO
{content}

## DIRETRIZES DE CRIAÇÃO
- Tom de voz: Técnico mas acessível, com foco no produtor rural
- Linguagem: Clara, objetiva, baseada em dados e ciência
- Diferencial: Destacar os benefícios práticos para o dia a dia do produtor

## PÚBLICO-ALVO
Produtores rurais, agrônomos, consultores e profissionais do agronegócio

## CALL TO ACTION
- Para mais informações, consulte seu representante técnico
- Visite nosso site para saber mais sobre {product_name}

## CONSIDERAÇÕES IMPORTANTES
- Baseie-se em dados técnicos e científicos
- Destaque os resultados comprovados em campo
- Mantenha o foco nos benefícios para o produtor

---
*Briefing gerado automaticamente pelo Agente Seeds*
"""
    return briefing

def is_syn_agent(agent_name):
    return agent_name and any(keyword in agent_name.upper() for keyword in ['SYN'])

def selecionar_agente_interface():
    st.title("Agente Seeds")
    
    agentes = listar_agentes()
    
    if not agentes:
        st.error("❌ Nenhum agente disponível. Crie um agente primeiro na aba de Gerenciamento.")
        return None
    
    opcoes_agentes = []
    for agente in agentes:
        agente_completo = obter_agente_com_heranca(agente['_id'])
        if agente_completo:
            descricao = f"{agente['nome']} - {agente.get('categoria', 'Social')}"
            if agente.get('agente_mae_id'):
                descricao += " 🔗"
            squad_permitido = agente.get('squad_permitido', 'Todos')
            descricao += f" 👥{squad_permitido}"
            opcoes_agentes.append((descricao, agente_completo))
    
    if opcoes_agentes:
        agente_selecionado_desc = st.selectbox(
            "Selecione uma base de conhecimento para usar o sistema:",
            options=[op[0] for op in opcoes_agentes],
            index=0,
            key="selectbox_agente_principal"
        )
        
        agente_completo = None
        for desc, agente in opcoes_agentes:
            if desc == agente_selecionado_desc:
                agente_completo = agente
                break
        
        if agente_completo and st.button("✅ Confirmar Seleção", key="confirmar_agente"):
            st.session_state.agente_selecionado = agente_completo
            st.session_state.messages = []
            st.session_state.segmentos_selecionados = ["system_prompt", "base_conhecimento", "comments", "planejamento"]
            st.success(f"✅ Agente '{agente_completo['nome']}' selecionado!")
            st.rerun()
        
        return agente_completo
    else:
        st.info("Nenhum agente disponível com as permissões atuais.")
        return None

# --- VERIFICAÇÃO DE AGENTE SELECIONADO ---
if not st.session_state.agente_selecionado:
    selecionar_agente_interface()
    st.stop()

# --- INTERFACE PRINCIPAL ---
agente_selecionado = st.session_state.agente_selecionado

# Sidebar com informações do usuário
st.sidebar.title(f"🤖 Bem-vindo, {get_current_user().get('nome', 'Usuário')}!")
st.sidebar.info(f"**Squad:** {get_current_squad()}")
st.sidebar.info(f"**Agente selecionado:** {agente_selecionado['nome']}")

if st.sidebar.button("🚪 Sair", key="logout_btn"):
    for key in ["logged_in", "user", "admin_password_correct", "admin_user", "agente_selecionado", "messages"]:
        if key in st.session_state:
            del st.session_state[key]
    st.rerun()

if st.sidebar.button("🔄 Trocar Agente", key="trocar_agente_global"):
    st.session_state.agente_selecionado = None
    st.session_state.messages = []
    st.rerun()

# --- STATUS DAS APIS NA SIDEBAR ---
st.sidebar.markdown("---")
st.sidebar.subheader("🔌 Status das APIs")

if gemini_api_key:
    st.sidebar.success("✅ Gemini: OK")
else:
    st.sidebar.error("❌ Gemini: Não configurado")

if anthropic_api_key:
    st.sidebar.success("✅ Claude: OK")
else:
    st.sidebar.error("❌ Claude: Não configurado")

if openai_api_key:
    st.sidebar.success("✅ OpenAI: OK")
else:
    st.sidebar.warning("⚠️ OpenAI: Não configurado")

if perp_api_key:
    st.sidebar.success("✅ Perplexity: OK")
else:
    st.sidebar.warning("⚠️ Perplexity: Não configurado")

# --- TÍTULO PRINCIPAL E SELECTOR DE AGENTE ---
st.title("🤖 Agente Seeds")

# Carregar agentes disponíveis
agentes = listar_agentes()

if agentes:
    opcoes_agentes = []
    for agente in agentes:
        agente_completo = obter_agente_com_heranca(agente['_id'])
        if agente_completo:
            descricao = f"{agente['nome']} - {agente.get('categoria', 'Social')}"
            if agente.get('agente_mae_id'):
                descricao += " 🔗"
            squad_permitido = agente.get('squad_permitido', 'Todos')
            descricao += f" 👥{squad_permitido}"
            opcoes_agentes.append((descricao, agente_completo))
    
    if opcoes_agentes:
        indice_atual = 0
        for i, (desc, agente) in enumerate(opcoes_agentes):
            if agente['_id'] == st.session_state.agente_selecionado['_id']:
                indice_atual = i
                break
        
        col1, col2 = st.columns([3, 1])
        with col1:
            novo_agente_desc = st.selectbox(
                "Selecionar Agente:",
                options=[op[0] for op in opcoes_agentes],
                index=indice_atual,
                key="selectbox_trocar_agente"
            )
        with col2:
            if st.button("🔄 Trocar", key="botao_trocar_agente"):
                for desc, agente in opcoes_agentes:
                    if desc == novo_agente_desc:
                        st.session_state.agente_selecionado = agente
                        st.session_state.messages = []
                        st.success(f"✅ Agente alterado para '{agente['nome']}'!")
                        st.rerun()
                        break

# --- DEFINIÇÃO DAS ABAS ---
abas_base = [
    "💬 Chat", 
    "⚙️ Gerenciar Agentes",
    "📓 Diário de Bordo",
    "✅ Validação Unificada",
    "✨ Geração de Conteúdo",
    "📝 Revisão Ortográfica",
    "Monitoramento de Redes",
    "🚀 Otimização de Conteúdo",
    "📅 Criadora de Calendário",
    "📊 Planejamento Estratégico",
    "📱 Planejamento de Mídias",
]

if is_syn_agent(agente_selecionado['nome']):
    abas_base.append("📋 Briefing")

tabs = st.tabs(abas_base)
tab_mapping = {aba: tabs[i] for i, aba in enumerate(abas_base)}

# ==================== ABA: CHAT ====================
with tab_mapping["💬 Chat"]:
    st.header("💬 Chat com Agente")
    
    if 'messages' not in st.session_state:
        st.session_state.messages = []
    if 'segmentos_selecionados' not in st.session_state:
        st.session_state.segmentos_selecionados = []
    if 'modelo_chat' not in st.session_state:
        st.session_state.modelo_chat = "Gemini"
    
    agente = st.session_state.agente_selecionado
    st.subheader(f"Conversando com: {agente['nome']}")
    
    st.sidebar.subheader("🤖 Configurações do Modelo")
    modelo_chat = st.sidebar.selectbox(
        "Escolha o modelo:",
        ["Gemini", "Claude", "OpenAI"],
        key="modelo_chat_selector",
        index=0 if st.session_state.modelo_chat == "Gemini" else 1 if st.session_state.modelo_chat == "Claude" else 2
    )
    st.session_state.modelo_chat = modelo_chat
    
    if modelo_chat == "Gemini" and not gemini_api_key:
        st.sidebar.error("❌ Gemini não disponível")
    elif modelo_chat == "Claude" and not anthropic_api_key:
        st.sidebar.error("❌ Claude não disponível")
    elif modelo_chat == "OpenAI" and not openai_api_key:
        st.sidebar.error("❌ OpenAI não disponível")
    else:
        st.sidebar.success(f"✅ {modelo_chat} ativo")
    
    st.sidebar.subheader("🔧 Configurações do Agente")
    st.sidebar.write("Selecione quais bases de conhecimento usar:")
    
    segmentos_disponiveis = {
        "Prompt do Sistema": "system_prompt",
        "Brand Guidelines": "base_conhecimento", 
        "Diário do Cliente": "comments",
        "Planejamento": "planejamento"
    }
    
    segmentos_selecionados = []
    for nome, chave in segmentos_disponiveis.items():
        if st.sidebar.checkbox(nome, value=chave in st.session_state.segmentos_selecionados, key=f"seg_{chave}"):
            segmentos_selecionados.append(chave)
    
    st.session_state.segmentos_selecionados = segmentos_selecionados
    
    if segmentos_selecionados:
        st.sidebar.success(f"✅ Usando {len(segmentos_selecionados)} segmento(s)")
    else:
        st.sidebar.warning("⚠️ Nenhum segmento selecionado")
    
    if hasattr(st.session_state, 'messages') and st.session_state.messages:
        for message in st.session_state.messages:
            if isinstance(message, dict) and "role" in message:
                with st.chat_message(message["role"]):
                    st.markdown(message.get("content", ""))
    else:
        st.info("💬 Inicie uma conversa digitando uma mensagem abaixo!")
    
    if prompt := st.chat_input("Digite sua mensagem..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        
        contexto = construir_contexto(
            agente, 
            st.session_state.segmentos_selecionados, 
            st.session_state.messages
        )
        
        with st.chat_message("assistant"):
            with st.spinner('Pensando...'):
                try:
                    resposta = gerar_resposta_modelo(
                        prompt, 
                        st.session_state.modelo_chat,
                        contexto
                    )
                    st.markdown(resposta)
                    
                    st.session_state.messages.append({"role": "assistant", "content": resposta})
                    
                    salvar_conversa(
                        agente['_id'], 
                        st.session_state.messages,
                        st.session_state.segmentos_selecionados
                    )
                    
                except Exception as e:
                    st.error(f"Erro ao gerar resposta: {str(e)}")

# ==================== ABA: GERENCIAR AGENTES ====================
with tab_mapping["⚙️ Gerenciar Agentes"]:
    st.header("Gerenciamento de Agentes")
    
    current_user = get_current_user()
    current_squad = get_current_squad()
    
    if current_squad not in ["admin", "Syngenta", "SME", "Enterprise"]:
        st.warning("Acesso restrito a usuários autorizados")
    else:
        if current_squad == "admin":
            if not check_admin_password():
                st.warning("Acesso restrito a administradores")
            else:
                st.write(f'Bem-vindo administrador!')
        else:
            st.write(f'Bem-vindo {current_user.get("nome", "Usuário")} do squad {current_squad}!')
        
        sub_tab1, sub_tab2, sub_tab3 = st.tabs(["Criar Agente", "Editar Agente", "Gerenciar Agentes"])
        
        with sub_tab1:
            st.subheader("Criar Novo Agente")
            
            with st.form("form_criar_agente"):
                nome_agente = st.text_input("Nome do Agente:")
                
                categoria = st.selectbox(
                    "Categoria:",
                    ["Social", "SEO", "Conteúdo", "Monitoramento"],
                    help="Organize o agente por área de atuação"
                )
                
                squad_permitido = st.selectbox(
                    "Squad Permitido:",
                    ["Todos", "Syngenta", "SME", "Enterprise"],
                    help="Selecione qual squad pode ver e usar este agente"
                )
                
                if categoria == "Monitoramento":
                    st.info("🔍 **Agente de Monitoramento**: Este agente será usado apenas na aba de Monitoramento de Redes.")
                    
                    base_conhecimento = st.text_area(
                        "Base de Conhecimento para Monitoramento:", 
                        height=300,
                        placeholder="Cole aqui a base de conhecimento específica para monitoramento de redes sociais."
                    )
                    
                    system_prompt = ""
                    comments = ""
                    planejamento = ""
                    criar_como_filho = False
                    agente_mae_id = None
                    herdar_elementos = []
                    
                else:
                    criar_como_filho = st.checkbox("Criar como agente filho (herdar elementos)")
                    
                    agente_mae_id = None
                    herdar_elementos = []
                    
                    if criar_como_filho:
                        agentes_mae = listar_agentes_para_heranca()
                        agentes_mae = [agente for agente in agentes_mae if agente.get('categoria') != 'Monitoramento']
                        
                        if agentes_mae:
                            agente_mae_options = {f"{agente['nome']} ({agente.get('categoria', 'Social')})": agente['_id'] for agente in agentes_mae}
                            agente_mae_selecionado = st.selectbox(
                                "Agente Mãe:",
                                list(agente_mae_options.keys()),
                                help="Selecione o agente do qual este agente irá herdar elementos"
                            )
                            agente_mae_id = agente_mae_options[agente_mae_selecionado]
                            
                            herdar_elementos = st.multiselect(
                                "Selecione os elementos a herdar do agente mãe:",
                                ["system_prompt", "base_conhecimento", "comments", "planejamento"],
                                help="Estes elementos serão herdados do agente mãe se não preenchidos abaixo"
                            )
                        else:
                            st.info("Nenhum agente disponível para herança. Crie primeiro um agente mãe.")
                    
                    system_prompt = st.text_area("Prompt de Sistema:", height=150)
                    base_conhecimento = st.text_area("Brand Guidelines:", height=200)
                    comments = st.text_area("Diário do cliente:", height=200)
                    planejamento = st.text_area("Planejamento:", height=200)
                
                submitted = st.form_submit_button("Criar Agente")
                if submitted:
                    if nome_agente:
                        agente_id = criar_agente(
                            nome_agente, 
                            system_prompt, 
                            base_conhecimento, 
                            comments, 
                            planejamento,
                            categoria,
                            squad_permitido,
                            agente_mae_id if criar_como_filho else None,
                            herdar_elementos if criar_como_filho else []
                        )
                        st.success(f"Agente '{nome_agente}' criado com sucesso!")
                    else:
                        st.error("Nome é obrigatório!")
        
        with sub_tab2:
            st.subheader("Editar Agente Existente")
            
            agentes = listar_agentes()
            if agentes:
                agente_options = {agente['nome']: agente for agente in agentes}
                agente_selecionado_nome = st.selectbox("Selecione o agente para editar:", 
                                                     list(agente_options.keys()))
                
                if agente_selecionado_nome:
                    agente = agente_options[agente_selecionado_nome]
                    
                    with st.form("form_editar_agente"):
                        novo_nome = st.text_input("Nome do Agente:", value=agente['nome'])
                        
                        categorias_disponiveis = ["Social", "SEO", "Conteúdo", "Monitoramento"]
                        index_categoria = categorias_disponiveis.index(agente.get('categoria', 'Social')) if agente.get('categoria') in categorias_disponiveis else 0
                        nova_categoria = st.selectbox("Categoria:", categorias_disponiveis, index=index_categoria)
                        
                        squads_disponiveis = ["Todos", "Syngenta", "SME", "Enterprise"]
                        squad_atual = agente.get('squad_permitido', 'Todos')
                        index_squad = squads_disponiveis.index(squad_atual) if squad_atual in squads_disponiveis else 0
                        novo_squad_permitido = st.selectbox("Squad Permitido:", squads_disponiveis, index=index_squad)
                        
                        if nova_categoria == "Monitoramento":
                            nova_base = st.text_area("Base de Conhecimento:", value=agente.get('base_conhecimento', ''), height=300)
                            novo_prompt = ""
                            nova_comment = ""
                            novo_planejamento = ""
                            agente_mae_id = None
                            herdar_elementos = []
                        else:
                            if agente.get('agente_mae_id'):
                                agente_mae = obter_agente(agente['agente_mae_id'])
                                if agente_mae:
                                    st.info(f"🔗 Este agente é filho de: {agente_mae['nome']}")
                            
                            novo_prompt = st.text_area("Prompt de Sistema:", value=agente.get('system_prompt', ''), height=150)
                            nova_base = st.text_area("Brand Guidelines:", value=agente.get('base_conhecimento', ''), height=200)
                            nova_comment = st.text_area("Diário:", value=agente.get('comments', ''), height=200)
                            novo_planejamento = st.text_area("Planejamento:", value=agente.get('planejamento', ''), height=200)
                            agente_mae_id = agente.get('agente_mae_id')
                            herdar_elementos = agente.get('herdar_elementos', [])
                        
                        submitted = st.form_submit_button("Atualizar Agente")
                        if submitted and novo_nome:
                            atualizar_agente(
                                agente['_id'], novo_nome, novo_prompt, nova_base, 
                                nova_comment, novo_planejamento, nova_categoria, novo_squad_permitido,
                                agente_mae_id, herdar_elementos
                            )
                            st.success(f"Agente '{novo_nome}' atualizado!")
                            st.rerun()
        
        with sub_tab3:
            st.subheader("Gerenciar Agentes")
            
            categorias = ["Todos", "Social", "SEO", "Conteúdo", "Monitoramento"]
            categoria_filtro = st.selectbox("Filtrar por categoria:", categorias)
            
            agentes = listar_agentes()
            if categoria_filtro != "Todos":
                agentes = [agente for agente in agentes if agente.get('categoria') == categoria_filtro]
            
            if agentes:
                for i, agente in enumerate(agentes):
                    with st.expander(f"{agente['nome']} - {agente.get('categoria', 'Social')} - Squad: {agente.get('squad_permitido', 'Todos')}"):
                        if current_squad == "admin" and agente.get('criado_por'):
                            st.write(f"**Proprietário:** {agente['criado_por']}")
                        
                        if agente.get('categoria') == 'Monitoramento':
                            st.info("🔍 Agente de Monitoramento")
                            if agente.get('base_conhecimento'):
                                st.write(f"**Base:** {agente['base_conhecimento'][:200]}...")
                        else:
                            if agente.get('agente_mae_id'):
                                st.write(f"**🔗 Herda de:** {agente.get('agente_mae_id')}")
                            st.write(f"**Prompt:** {agente.get('system_prompt', '')[:100]}...")
                        
                        col1, col2 = st.columns(2)
                        with col1:
                            if st.button("Selecionar para Chat", key=f"select_{i}"):
                                agente_completo = obter_agente_com_heranca(agente['_id'])
                                st.session_state.agente_selecionado = agente_completo
                                st.session_state.messages = []
                                st.success(f"Agente '{agente['nome']}' selecionado!")
                                st.rerun()
                        with col2:
                            if st.button("Desativar", key=f"delete_{i}"):
                                desativar_agente(agente['_id'])
                                st.success(f"Agente '{agente['nome']}' desativado!")
                                st.rerun()
            else:
                st.info("Nenhum agente encontrado.")

# ==================== ABA: DIÁRIO DE BORDO ====================
with tab_mapping["📓 Diário de Bordo"]:
    st.header("📓 Diário de Bordo - Cliente")
    
    agente = st.session_state.agente_selecionado
    st.subheader(f"Diário para: {agente['nome']}")
    
    comentarios_atuais = agente.get('comments', '')
    
    tab_visualizar, tab_adicionar, tab_relatorio = st.tabs(["👁️ Visualizar", "➕ Adicionar", "📊 Relatório"])
    
    with tab_visualizar:
        if comentarios_atuais:
            st.markdown("### 📝 Diário Atual do Cliente")
            
            palavras = len(comentarios_atuais.split())
            caracteres = len(comentarios_atuais)
            
            col_stat1, col_stat2 = st.columns(2)
            with col_stat1:
                st.metric("📝 Palavras", palavras)
            with col_stat2:
                st.metric("🔤 Caracteres", caracteres)
            
            st.text_area("Conteúdo do diário:", value=comentarios_atuais, height=400, disabled=True)
            
            st.download_button(
                "💾 Exportar Diário",
                data=comentarios_atuais,
                file_name=f"diario_{agente['nome']}_{datetime.datetime.now().strftime('%Y%m%d')}.txt",
                mime="text/plain"
            )
        else:
            st.info("📭 O diário está vazio.")
    
    with tab_adicionar:
        st.markdown("### 📤 Adicionar ao Diário")
        
        metodo_adicional = st.radio(
            "Como deseja adicionar conteúdo:",
            ["📝 Texto Manual", "📎 Upload de Documento", "✂️ Extrair de Conversa"],
            horizontal=True
        )
        
        if metodo_adicional == "📝 Texto Manual":
            data_registro = st.date_input("Data do registro:", value=datetime.datetime.now())
            titulo_registro = st.text_input("Título/Contexto:")
            
            novo_conteudo = st.text_area("Conteúdo:", height=200)
            
            if st.button("💾 Salvar no Diário", type="primary"):
                if novo_conteudo.strip():
                    entrada_formatada = f"\n\n--- {titulo_registro if titulo_registro else 'Nova Entrada'} ({data_registro.strftime('%d/%m/%Y')}) ---\n{novo_conteudo}"
                    novos_comentarios = comentarios_atuais + entrada_formatada
                    
                    atualizar_agente(
                        agente['_id'], agente['nome'], agente.get('system_prompt', ''),
                        agente.get('base_conhecimento', ''), novos_comentarios,
                        agente.get('planejamento', ''), agente.get('categoria', 'Social'),
                        agente.get('squad_permitido', 'Todos'), agente.get('agente_mae_id'),
                        agente.get('herdar_elementos', [])
                    )
                    st.session_state.agente_selecionado = obter_agente_com_heranca(agente['_id'])
                    st.success("✅ Conteúdo adicionado ao diário!")
                    st.rerun()
                else:
                    st.warning("Digite algum conteúdo para salvar")
        
        elif metodo_adicional == "📎 Upload de Documento":
            uploaded_file = st.file_uploader("Selecione um documento:", type=['pdf', 'docx', 'txt'])
            
            if uploaded_file:
                st.success(f"✅ {uploaded_file.name} carregado")
                
                with st.spinner("Extraindo texto..."):
                    try:
                        if uploaded_file.type == "application/pdf":
                            texto_extraido, _ = extract_text_from_pdf_com_slides(uploaded_file)
                        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                            texto_extraido = extrair_texto_arquivo(uploaded_file)
                        elif uploaded_file.type == "text/plain":
                            texto_extraido = str(uploaded_file.read(), "utf-8")
                        else:
                            texto_extraido = ""
                        
                        contexto_doc = st.text_input("Contexto/Origem:")
                        
                        if st.button("💾 Adicionar Documento ao Diário", type="primary"):
                            if texto_extraido.strip():
                                data_atual = datetime.datetime.now().strftime('%d/%m/%Y')
                                entrada_formatada = f"\n\n--- {contexto_doc if contexto_doc else 'Documento'} - {uploaded_file.name} ({data_atual}) ---\n{texto_extraido[:10000]}"
                                novos_comentarios = comentarios_atuais + entrada_formatada
                                
                                atualizar_agente(
                                    agente['_id'], agente['nome'], agente.get('system_prompt', ''),
                                    agente.get('base_conhecimento', ''), novos_comentarios,
                                    agente.get('planejamento', ''), agente.get('categoria', 'Social'),
                                    agente.get('squad_permitido', 'Todos'), agente.get('agente_mae_id'),
                                    agente.get('herdar_elementos', [])
                                )
                                st.session_state.agente_selecionado = obter_agente_com_heranca(agente['_id'])
                                st.success("✅ Documento adicionado ao diário!")
                                st.rerun()
                            else:
                                st.warning("Documento vazio ou não foi possível extrair texto")
                    except Exception as e:
                        st.error(f"Erro: {str(e)}")
        
        elif metodo_adicional == "✂️ Extrair de Conversa":
            conversas = obter_conversas(agente['_id'], limite=5)
            
            if conversas:
                for i, conversa in enumerate(conversas):
                    with st.expander(f"Conversa {i+1} - {conversa.get('data_criacao', 'Data desconhecida')}", expanded=False):
                        mensagens = conversa.get('mensagens', [])
                        for msg in mensagens[-6:]:
                            role = "👤" if msg.get("role") == "user" else "🤖"
                            st.write(f"{role}: {msg.get('content', '')[:200]}...")
                        
                        if st.button(f"📋 Usar esta conversa", key=f"usar_conversa_{i}"):
                            texto_conversa = ""
                            for msg in mensagens:
                                if msg.get("role") == "user":
                                    texto_conversa += f"Cliente: {msg.get('content', '')}\n"
                            
                            if texto_conversa.strip():
                                data_atual = datetime.datetime.now().strftime('%d/%m/%Y')
                                entrada_formatada = f"\n\n--- Conversa extraída ({data_atual}) ---\n{texto_conversa}"
                                novos_comentarios = comentarios_atuais + entrada_formatada
                                
                                atualizar_agente(
                                    agente['_id'], agente['nome'], agente.get('system_prompt', ''),
                                    agente.get('base_conhecimento', ''), novos_comentarios,
                                    agente.get('planejamento', ''), agente.get('categoria', 'Social'),
                                    agente.get('squad_permitido', 'Todos'), agente.get('agente_mae_id'),
                                    agente.get('herdar_elementos', [])
                                )
                                st.session_state.agente_selecionado = obter_agente_com_heranca(agente['_id'])
                                st.success("✅ Conversa adicionada ao diário!")
                                st.rerun()
            else:
                st.info("Nenhuma conversa recente encontrada")
    
    with tab_relatorio:
        st.markdown("### 📊 Relatório de Andamento com Cliente")
        
        if not comentarios_atuais or len(comentarios_atuais.strip()) < 50:
            st.info("📭 Diário muito curto para gerar relatório.")
        else:
            tipo_analise = st.selectbox("Tipo de análise:", ["Análise Completa", "Foco em Oportunidades", "Identificar Problemas"])
            formato_relatorio = st.selectbox("Formato do relatório:", ["Relatório Executivo", "Lista de Ações", "Análise Detalhada"])
            
            if st.button("📈 Gerar Análise do Diário", type="primary"):
                with st.spinner("Analisando diário..."):
                    try:
                        prompt_analise = f"""
                        ## ANÁLISE DE DIÁRIO DE CLIENTE
                        
                        **AGENTE:** {agente['nome']}
                        **TIPO DE ANÁLISE:** {tipo_analise}
                        **FORMATO:** {formato_relatorio}
                        
                        **CONTEÚDO DO DIÁRIO:**
                        {comentarios_atuais[:8000]}
                        
                        ## INSTRUÇÕES:
                        Analise o diário e gere um relatório identificando:
                        1. Padrões e tendências no feedback
                        2. Oportunidades de melhoria
                        3. Pontos críticos
                        4. Recomendações para próximos passos
                        """
                        
                        resposta = modelo_texto.generate_content(prompt_analise)
                        st.markdown(resposta.text)
                        
                        st.download_button(
                            "💾 Baixar Relatório",
                            data=resposta.text,
                            file_name=f"analise_diario_{agente['nome']}_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                            mime="text/plain"
                        )
                    except Exception as e:
                        st.error(f"Erro: {str(e)}")

# ==================== ABA: VALIDAÇÃO UNIFICADA ====================
with tab_mapping["✅ Validação Unificada"]:
    st.header("✅ Validação Unificada de Conteúdo")
    
    agente = st.session_state.agente_selecionado
    st.subheader(f"Validação com: {agente.get('nome', 'Agente')}")
    
    st.markdown("---")
    st.subheader("🎯 Contexto para Análise")
    
    contexto_global = st.text_area(
        "**✍️ Contexto adicional para todas as análises:**", 
        height=120, 
        key="contexto_global_validacao",
        placeholder="Forneça contexto adicional que será aplicado a TODAS as análises..."
    )
    
    subtab_imagem, subtab_texto, subtab_video, subtab_texto_imagem = st.tabs(
        ["🖼️ Validação de Imagem", "📄 Validação de Documentos", "🎬 Validação de Vídeo", "📝 Validação de Texto em Imagem"]
    )
    
    # SUBTAB: VALIDAÇÃO DE TEXTO EM IMAGEM
    with subtab_texto_imagem:
        st.subheader("📝 Validação de Texto em Imagem")
        
        uploaded_images_texto = st.file_uploader(
            "Carregue imagens para análise de texto",
            type=["jpg", "jpeg", "png", "webp"],
            accept_multiple_files=True,
            key="image_text_upload"
        )
        
        if st.button("🗑️ Limpar Análises", key="limpar_texto_imagem"):
            if 'resultados_texto_imagem' in st.session_state:
                del st.session_state.resultados_texto_imagem
            st.rerun()
        
        if uploaded_images_texto:
            st.success(f"✅ {len(uploaded_images_texto)} imagem(ns) carregada(s)")
            
            if st.button("🔍 Validar Texto em Todas as Imagens", type="primary"):
                resultados = []
                progress_bar = st.progress(0)
                
                for idx, uploaded_image in enumerate(uploaded_images_texto):
                    progress_bar.progress((idx + 1) / len(uploaded_images_texto))
                    
                    with st.spinner(f'Processando imagem {idx+1}...'):
                        try:
                            prompt_texto_imagem = f"""
                            {contexto_global if contexto_global else ''}
                            
                            ## ANÁLISE DE TEXTO EM IMAGEM
                            
                            **INSTRUÇÕES:**
                            1. Transcreva e analise TODO o texto visível na imagem
                            2. Foque em: ortografia, gramática, clareza
                            
                            **FORMATO DE RESPOSTA:**
                            
                            ## Arte {idx+1}
                            
                            **Texto:** "[Texto extraído]"
                            
                            **Correções:** [✅/⚠️/❌] [Descrição]
                            """
                            
                            response = modelo_vision.generate_content([
                                prompt_texto_imagem,
                                {"mime_type": uploaded_image.type, "data": uploaded_image.getvalue()}
                            ])
                            
                            analise = response.text
                            
                            if "❌" in analise:
                                status = "Com erros"
                            elif "⚠️" in analise:
                                status = "Ajustes sugeridos"
                            else:
                                status = "Correto"
                            
                            resultados.append({
                                'indice': idx + 1,
                                'nome': uploaded_image.name,
                                'analise': analise,
                                'status': status,
                                'imagem': uploaded_image
                            })
                            
                        except Exception as e:
                            resultados.append({
                                'indice': idx + 1,
                                'nome': uploaded_image.name,
                                'analise': f"❌ Erro: {str(e)}",
                                'status': "Erro",
                                'imagem': uploaded_image
                            })
                
                progress_bar.empty()
                st.session_state.resultados_texto_imagem = resultados
                
                relatorio_consolidado = gerar_relatorio_texto_imagem_consolidado(resultados)
                
                for resultado in resultados:
                    with st.expander(f"🖼️ Arte {resultado['indice']} - {resultado['status']}", expanded=True):
                        col_img, col_text = st.columns([1, 2])
                        with col_img:
                            image = Image.open(resultado['imagem'])
                            st.image(image, use_container_width=True)
                        with col_text:
                            st.markdown(resultado['analise'])
                
                st.download_button(
                    "📥 Baixar Relatório",
                    data=relatorio_consolidado,
                    file_name=f"relatorio_texto_imagens_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                    mime="text/plain"
                )
    
    # SUBTAB: VALIDAÇÃO DE DOCUMENTOS
    with subtab_texto:
        st.subheader("📄 Validação de Documentos e Texto")
        
        with st.expander("Configurações de Exportação PDF", expanded=True):
            incluir_comentarios_pdf = st.checkbox("Incluir comentários no PDF", value=True)
            gerar_relatorio_completo = st.checkbox("Gerar relatório completo", value=True)
            limitar_comentarios = st.slider("Máximo de comentários por PDF:", 1, 10, 5)
        
        if st.button("🗑️ Limpar Análises de Texto", key="limpar_analises_texto"):
            st.session_state.validacao_triggered = False
            st.session_state.todos_textos = []
            st.rerun()
        
        col_entrada, col_saida = st.columns([1, 1])
        
        with col_entrada:
            st.markdown("### Entrada de Conteúdo")
            texto_input = st.text_area("**Digite o texto para validação:**", height=150, key="texto_validacao")
            
            st.markdown("### 📎 Ou carregue arquivos")
            arquivos_documentos = st.file_uploader(
                "Documentos suportados: PDF, PPTX, TXT, DOCX",
                type=['pdf', 'pptx', 'txt', 'docx'],
                accept_multiple_files=True,
                key="arquivos_documentos_validacao"
            )
            
            with st.expander("Configurações de Análise de Texto"):
                analise_especializada = st.checkbox("Análise especializada por áreas", value=st.session_state.analise_especializada_texto)
                analisadores_selecionados = st.multiselect(
                    "Especialistas de texto a incluir:",
                    options=['ortografia', 'lexico', 'branding'],
                    default=st.session_state.analisadores_selecionados_texto,
                    format_func=lambda x: {'ortografia': 'Ortografia', 'lexico': 'Léxico', 'branding': 'Branding'}[x]
                )
            
            if st.button("Validar Conteúdo de Texto", type="primary", key="validate_documents", use_container_width=True):
                st.session_state.validacao_triggered = True
                st.session_state.analise_especializada_texto = analise_especializada
                st.session_state.analisadores_selecionados_texto = analisadores_selecionados
        
        with col_saida:
            st.markdown("### 📊 Resultados de Texto")
            
            if st.session_state.validacao_triggered:
                todos_textos = []
                resultados_pdf = {}
                
                if texto_input and texto_input.strip():
                    todos_textos.append({'nome': 'Texto_Manual', 'conteudo': texto_input, 'tipo': 'texto_direto'})
                
                if arquivos_documentos:
                    for arquivo in arquivos_documentos:
                        with st.spinner(f"Processando {arquivo.name}..."):
                            try:
                                if arquivo.type == "application/pdf":
                                    texto_extraido, slides_info = extract_text_from_pdf_com_slides(arquivo)
                                    arquivo_original = arquivo
                                elif arquivo.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                                    texto_extraido, slides_info = extract_text_from_pptx_com_slides(arquivo)
                                    arquivo_original = None
                                else:
                                    texto_extraido = extrair_texto_arquivo(arquivo)
                                    slides_info = []
                                    arquivo_original = None
                                
                                if texto_extraido and texto_extraido.strip():
                                    todos_textos.append({
                                        'nome': arquivo.name,
                                        'conteudo': texto_extraido,
                                        'slides': slides_info,
                                        'tipo': arquivo.type,
                                        'arquivo_original': arquivo_original
                                    })
                            except Exception as e:
                                st.error(f"Erro: {str(e)}")
                
                if not todos_textos:
                    st.warning("Nenhum conteúdo válido encontrado.")
                else:
                    st.success(f"{len(todos_textos)} documento(s) processado(s)")
                    
                    for doc in todos_textos:
                        with st.expander(f"📄 {doc['nome']}", expanded=True):
                            contexto_agente = ""
                            if "base_conhecimento" in agente:
                                contexto_agente = agente['base_conhecimento']
                            
                            contexto_completo = contexto_agente
                            if contexto_global:
                                contexto_completo += f"\n\n{contexto_global}"
                            
                            with st.spinner(f"Analisando {doc['nome']}..."):
                                try:
                                    if st.session_state.analise_especializada_texto:
                                        analisadores_config = criar_analisadores_texto(contexto_completo, "")
                                        analisadores_filtrados = {k: v for k, v in analisadores_config.items() 
                                                                 if k in st.session_state.analisadores_selecionados_texto}
                                        
                                        resultados_especialistas = executar_analise_texto_especializada(
                                            doc['conteudo'], doc['nome'], analisadores_filtrados
                                        )
                                        
                                        relatorio = gerar_relatorio_texto_consolidado(resultados_especialistas, doc['nome'])
                                        st.markdown(relatorio)
                                        
                                        if incluir_comentarios_pdf and doc.get('arquivo_original'):
                                            comentarios = extrair_comentarios_analise(relatorio)
                                            if comentarios:
                                                pdf_com_comentarios = adicionar_comentarios_pdf(
                                                    doc['arquivo_original'], comentarios[:limitar_comentarios], doc['nome']
                                                )
                                                if pdf_com_comentarios:
                                                    st.download_button(
                                                        "Baixar PDF com Comentários",
                                                        data=pdf_com_comentarios.getvalue(),
                                                        file_name=f"comentarios_{doc['nome']}",
                                                        mime="application/pdf",
                                                        key=f"download_pdf_{doc['nome']}"
                                                    )
                                    else:
                                        prompt = f"""
                                        {contexto_completo}
                                        
                                        Analise o texto e forneça validação.
                                        
                                        TEXTO:
                                        {doc['conteudo'][:8000]}
                                        
                                        FORMATO:
                                        ### CONFORMIDADE
                                        ### INCONSISTÊNCIAS
                                        ### SUGESTÕES
                                        """
                                        resposta = modelo_texto.generate_content(prompt)
                                        st.markdown(resposta.text)
                                        
                                except Exception as e:
                                    st.error(f"Erro: {str(e)}")
    
    # SUBTAB: VALIDAÇÃO DE IMAGEM
    with subtab_imagem:
        st.subheader("🖼️ Validação de Imagem")
        
        if st.button("🗑️ Limpar Análises de Imagem", key="limpar_analises_imagem"):
            st.session_state.resultados_analise_imagem = []
            st.rerun()
        
        uploaded_images = st.file_uploader(
            "Carregue imagens para análise", 
            type=["jpg", "jpeg", "png", "webp"], 
            key="image_upload_validacao",
            accept_multiple_files=True
        )
        
        with st.expander("⚙️ Configurações de Análise de Imagem"):
            analise_especializada_imagem = st.checkbox("Análise especializada", value=st.session_state.analise_especializada_imagem)
            analisadores_selecionados_imagem = st.multiselect(
                "Especialistas de imagem:",
                options=['composicao_visual', 'cores_branding'],
                default=st.session_state.analisadores_selecionados_imagem,
                format_func=lambda x: {'composicao_visual': 'Composição Visual', 'cores_branding': 'Cores e Branding'}[x]
            )
        
        if uploaded_images:
            if st.button("🔍 Validar Todas as Imagens", type="primary"):
                resultados_analise = []
                
                for idx, uploaded_image in enumerate(uploaded_images):
                    with st.spinner(f'Analisando imagem {idx+1}...'):
                        try:
                            with st.container():
                                st.markdown("---")
                                col_img, col_info = st.columns([2, 1])
                                
                                with col_img:
                                    image = Image.open(uploaded_image)
                                    st.image(image, use_container_width=True)
                                
                                with col_info:
                                    st.metric("📐 Dimensões", f"{image.width} x {image.height}")
                                
                                with st.expander(f"📋 Análise Detalhada", expanded=True):
                                    contexto_agente = ""
                                    if "base_conhecimento" in agente:
                                        contexto_agente = agente['base_conhecimento']
                                    
                                    contexto_completo = contexto_agente
                                    if contexto_global:
                                        contexto_completo += f"\n\n{contexto_global}"
                                    
                                    if st.session_state.analise_especializada_imagem:
                                        analisadores_config = criar_analisadores_imagem(contexto_completo, "")
                                        analisadores_filtrados = {k: v for k, v in analisadores_config.items() 
                                                                 if k in st.session_state.analisadores_selecionados_imagem}
                                        
                                        resultados_especialistas = executar_analise_imagem_especializada(
                                            uploaded_image, uploaded_image.name, analisadores_filtrados
                                        )
                                        
                                        relatorio = gerar_relatorio_imagem_consolidado(
                                            resultados_especialistas, uploaded_image.name, f"{image.width}x{image.height}"
                                        )
                                        st.markdown(relatorio)
                                        
                                        resultados_analise.append({
                                            'nome': uploaded_image.name,
                                            'analise': relatorio
                                        })
                                        
                        except Exception as e:
                            st.error(f"Erro: {str(e)}")
                
                st.session_state.resultados_analise_imagem = resultados_analise
    
    # SUBTAB: VALIDAÇÃO DE VÍDEO
    with subtab_video:
        st.subheader("🎬 Validação de Vídeo")
        
        if st.button("🗑️ Limpar Análises de Vídeo", key="limpar_analises_video"):
            st.session_state.resultados_analise_video = []
            st.rerun()
        
        col_upload, col_config = st.columns([2, 1])
        
        with col_upload:
            uploaded_videos = st.file_uploader(
                "Carregue vídeos para análise",
                type=["mp4", "mpeg", "mov", "avi", "webm"],
                key="video_upload_validacao",
                accept_multiple_files=True
            )
        
        with col_config:
            st.markdown("### ⚙️ Configurações")
            contexto_video_especifico = st.text_area("Contexto específico para vídeos:", height=100, key="video_context_especifico")
            analise_especializada_video = st.checkbox("Análise especializada", value=True)
            analisadores_selecionados_video = st.multiselect(
                "Especialistas de vídeo:",
                options=['narrativa_estrutura', 'qualidade_audio'],
                default=['narrativa_estrutura', 'qualidade_audio'],
                format_func=lambda x: {'narrativa_estrutura': 'Narrativa', 'qualidade_audio': 'Qualidade de Áudio'}[x]
            )
        
        if uploaded_videos:
            if st.button("🎬 Validar Todos os Vídeos", type="primary"):
                resultados_video = []
                
                for idx, uploaded_video in enumerate(uploaded_videos):
                    with st.spinner(f'Analisando vídeo {idx+1}...'):
                        try:
                            with st.container():
                                st.markdown("---")
                                st.subheader(f"🎬 {uploaded_video.name}")
                                
                                with st.expander("👀 Preview", expanded=False):
                                    st.video(uploaded_video)
                                
                                with st.expander(f"📋 Análise Completa", expanded=True):
                                    contexto_agente = ""
                                    if "base_conhecimento" in agente:
                                        contexto_agente = agente['base_conhecimento']
                                    
                                    contexto_completo = contexto_agente
                                    if contexto_global:
                                        contexto_completo += f"\n\n{contexto_global}"
                                    if contexto_video_especifico:
                                        contexto_completo += f"\n\n{contexto_video_especifico}"
                                    
                                    analisadores_config = criar_analisadores_video(contexto_agente, contexto_global, contexto_video_especifico)
                                    analisadores_filtrados = {k: v for k, v in analisadores_config.items() 
                                                             if k in analisadores_selecionados_video}
                                    
                                    resultados_especialistas = executar_analise_video_especializada(
                                        uploaded_video, uploaded_video.name, analisadores_filtrados
                                    )
                                    
                                    relatorio = gerar_relatorio_video_consolidado(
                                        resultados_especialistas, uploaded_video.name, uploaded_video.type
                                    )
                                    st.markdown(relatorio)
                                    
                                    resultados_video.append({'nome': uploaded_video.name, 'analise': relatorio})
                                    
                        except Exception as e:
                            st.error(f"Erro: {str(e)}")
                
                st.session_state.resultados_analise_video = resultados_video

# ==================== ABA: GERAÇÃO DE CONTEÚDO ====================
with tab_mapping["✨ Geração de Conteúdo"]:
    st.header("✨ Geração de Conteúdo com Múltiplos Insumos")
    
    def extrair_texto_arquivo_gen(arquivo):
        try:
            if arquivo.type == "text/plain":
                return str(arquivo.read(), "utf-8")
            elif arquivo.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                doc = docx.Document(io.BytesIO(arquivo.read()))
                return "\n".join([para.text for para in doc.paragraphs])
            else:
                return f"Tipo não suportado"
        except Exception as e:
            return f"Erro: {str(e)}"
    
    tab_geracao, tab_ajuste = st.tabs(["📝 Geração de Conteúdo", "✏️ Ajustes Incrementais"])
    
    with tab_geracao:
        col1, col2 = st.columns([2, 1])
        
        with col1:
            st.subheader("📝 Fontes de Conteúdo")
            
            usar_busca_web = st.checkbox("🔍 Realizar busca web", value=True)
            
            if usar_busca_web:
                termos_busca = st.text_area("🔎 Termos para busca web:", height=100, placeholder="Ex: tendências marketing digital 2024...")
            
            arquivos_upload = st.file_uploader(
                "Upload de Arquivos:",
                type=['pdf', 'txt', 'pptx', 'docx'],
                accept_multiple_files=True,
                key="arquivos_conteudo"
            )
            
            textos_arquivos = ""
            if arquivos_upload:
                for arquivo in arquivos_upload:
                    texto_extraido = extrair_texto_arquivo_gen(arquivo)
                    textos_arquivos += f"\n\n--- {arquivo.name} ---\n{texto_extraido}"
            
            briefing_manual = st.text_area("Briefing Manual:", height=150, key="briefing_manual")
        
        with col2:
            st.subheader("⚙️ Configurações")
            
            modelo_principal = st.selectbox("Modelo:", ["Gemini", "Claude", "OpenAI"], key="modelo_principal_select")
            
            modo_geracao = st.radio("Modo de Geração:", ["Configurações Padrão", "Prompt Personalizado"])
            
            if modo_geracao == "Configurações Padrão":
                tipo_conteudo = st.selectbox("Tipo de Conteúdo:", ["Post Social", "Artigo Blog", "Email Marketing", "Landing Page"])
                tom_voz = st.text_area("Tom de Voz:", placeholder="Ex: Formal e profissional")
                numero_palavras = st.slider("Número de Palavras:", 100, 3000, 800)
                usar_contexto_agente = st.checkbox("Usar contexto do agente", value=True)
            else:
                prompt_personalizado = st.text_area("Prompt Personalizado:", height=200)
        
        if st.button("🚀 Gerar Conteúdo", type="primary", use_container_width=True):
            tem_conteudo = (arquivos_upload or briefing_manual or (usar_busca_web and termos_busca))
            
            if not tem_conteudo:
                st.warning("Forneça pelo menos uma fonte de conteúdo")
            else:
                with st.spinner("Gerando conteúdo..."):
                    try:
                        contexto_completo = ""
                        
                        if textos_arquivos:
                            contexto_completo += f"### ARQUIVOS:\n{textos_arquivos}\n\n"
                        
                        if briefing_manual:
                            contexto_completo += f"### BRIEFING:\n{briefing_manual}\n\n"
                        
                        if usar_busca_web and termos_busca:
                            busca_resultado = realizar_busca_web_com_fontes(termos_busca, "")
                            contexto_completo += f"### BUSCA WEB:\n{busca_resultado}\n\n"
                        
                        contexto_agente = ""
                        if usar_contexto_agente and st.session_state.agente_selecionado:
                            contexto_agente = construir_contexto(agente_selecionado, st.session_state.segmentos_selecionados)
                        
                        if modo_geracao == "Configurações Padrão":
                            prompt = f"""
                            {contexto_agente}
                            
                            ## FONTES:
                            {contexto_completo}
                            
                            ## INSTRUÇÕES:
                            Tipo: {tipo_conteudo}
                            Tom: {tom_voz}
                            Palavras: {numero_palavras}
                            
                            Gere o conteúdo solicitado.
                            """
                        else:
                            prompt = f"""
                            {contexto_agente}
                            
                            {prompt_personalizado}
                            
                            ## FONTES:
                            {contexto_completo}
                            """
                        
                        conteudo = gerar_conteudo_modelo(prompt, modelo_principal, contexto_agente)
                        st.session_state.conteudo_gerado = conteudo
                        
                        st.subheader("📄 Conteúdo Gerado")
                        st.write(conteudo)
                        
                        st.download_button(
                            "💾 Baixar Conteúdo",
                            data=conteudo,
                            file_name=f"conteudo_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                            mime="text/plain"
                        )
                        
                    except Exception as e:
                        st.error(f"Erro: {str(e)}")
    
    with tab_ajuste:
        st.header("✏️ Ajustes Incrementais")
        
        if 'conteudo_gerado' not in st.session_state or not st.session_state.conteudo_gerado:
            st.info("Nenhum conteúdo gerado. Gere um conteúdo primeiro.")
        else:
            st.write(f"**Modelo:** {st.session_state.get('modelo_principal_select', 'Gemini')}")
            
            instrucoes_ajuste = st.text_area("Descreva os ajustes desejados:", height=100)
            modelo_ajuste = st.selectbox("Modelo para ajuste:", ["Gemini", "Claude", "OpenAI"])
            
            if st.button("🔄 Aplicar Ajustes", type="primary"):
                if instrucoes_ajuste:
                    with st.spinner("Aplicando ajustes..."):
                        try:
                            contexto_agente = ""
                            if st.session_state.agente_selecionado:
                                contexto_agente = construir_contexto(agente_selecionado, st.session_state.segmentos_selecionados)
                            
                            prompt = f"""
                            {contexto_agente}
                            
                            ## CONTEÚDO ORIGINAL:
                            {st.session_state.conteudo_gerado}
                            
                            ## AJUSTES SOLICITADOS:
                            {instrucoes_ajuste}
                            
                            Aplique os ajustes e retorne o conteúdo modificado.
                            """
                            
                            conteudo_ajustado = gerar_conteudo_modelo(prompt, modelo_ajuste, contexto_agente)
                            st.session_state.conteudo_gerado = conteudo_ajustado
                            
                            st.subheader("📄 Conteúdo Ajustado")
                            st.write(conteudo_ajustado)
                            
                            st.download_button(
                                "💾 Baixar Conteúdo Ajustado",
                                data=conteudo_ajustado,
                                file_name=f"conteudo_ajustado_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                                mime="text/plain"
                            )
                            
                        except Exception as e:
                            st.error(f"Erro: {str(e)}")
                else:
                    st.warning("Descreva os ajustes desejados")

# ==================== ABA: REVISÃO ORTOGRÁFICA ====================
with tab_mapping["📝 Revisão Ortográfica"]:
    st.header("📝 Revisão Ortográfica e Gramatical")
    
    modelo_revisao = st.sidebar.selectbox("Modelo:", ["Gemini", "Claude", "OpenAI"], key="modelo_revisao_selector")
    
    agente = st.session_state.agente_selecionado
    st.subheader(f"Revisão com: {agente['nome']}")
    
    st.sidebar.subheader("🔧 Configurações")
    segmentos_revisao = st.sidebar.multiselect(
        "Bases para revisão:",
        options=["system_prompt", "base_conhecimento", "comments", "planejamento"],
        default=st.session_state.get('segmentos_selecionados', [])
    )
    
    tab_texto, tab_arquivo = st.tabs(["📝 Texto Direto", "📎 Upload de Arquivos"])
    
    with tab_texto:
        col_original, col_resultado = st.columns(2)
        
        with col_original:
            st.subheader("📄 Texto Original")
            texto_para_revisao = st.text_area("Cole o texto:", height=400, key="texto_revisao")
            
            if texto_para_revisao:
                palavras = len(texto_para_revisao.split())
                st.metric("📊 Palavras", palavras)
            
            with st.expander("⚙️ Configurações"):
                revisao_estilo = st.checkbox("Incluir revisão de estilo", value=True)
                explicar_alteracoes = st.checkbox("Explicar alterações", value=True)
        
        with col_resultado:
            st.subheader("📋 Resultado da Revisão")
            
            if st.button("🔍 Realizar Revisão", type="primary"):
                if texto_para_revisao.strip():
                    with st.spinner("Revisando texto..."):
                        resultado = revisar_texto_ortografia(
                            texto_para_revisao, agente, segmentos_revisao,
                            revisao_estilo, True, explicar_alteracoes, modelo_revisao
                        )
                        st.markdown(resultado)
                        
                        st.download_button(
                            "💾 Baixar Relatório",
                            data=resultado,
                            file_name=f"revisao_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                            mime="text/plain"
                        )
                else:
                    st.warning("Cole o texto para revisão")
    
    with tab_arquivo:
        st.subheader("📎 Upload de Arquivos")
        
        arquivos_upload = st.file_uploader(
            "Selecione arquivos:", type=['pdf', 'pptx', 'txt', 'docx'],
            accept_multiple_files=True, key="arquivos_revisao"
        )
        
        if arquivos_upload:
            if st.button("🔍 Revisar Todos os Arquivos", type="primary"):
                for arquivo in arquivos_upload:
                    with st.spinner(f"Processando {arquivo.name}..."):
                        try:
                            if arquivo.type == "application/pdf":
                                texto_extraido, _ = extract_text_from_pdf_com_slides(arquivo)
                            elif arquivo.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                                texto_extraido, _ = extract_text_from_pptx_com_slides(arquivo)
                            else:
                                texto_extraido = extrair_texto_arquivo(arquivo)
                            
                            if texto_extraido and texto_extraido.strip():
                                with st.expander(f"📄 {arquivo.name}", expanded=False):
                                    resultado = revisar_texto_ortografia(
                                        texto_extraido, agente, segmentos_revisao,
                                        True, True, True, modelo_revisao
                                    )
                                    st.markdown(resultado)
                        except Exception as e:
                            st.error(f"Erro em {arquivo.name}: {str(e)}")

# ==================== ABA: MONITORAMENTO DE REDES ====================
with tab_mapping["Monitoramento de Redes"]:
    st.header("🤖 Agente de Monitoramento")
    st.markdown("**Especialista que fala como gente**")
    
    contexto_adicional = st.text_area(
        "📝 Contexto Adicional para Respostas:",
        height=100,
        placeholder="Ex: Este post é sobre vagas de emprego...",
        key="contexto_monitoramento"
    )
    
    modelo_monitoramento = st.sidebar.selectbox("Modelo:", ["Gemini", "Claude"], key="modelo_monitoramento_selector")
    
    agentes_monitoramento = [agente for agente in listar_agentes() if agente.get('categoria') == 'Monitoramento']
    
    if agentes_monitoramento:
        opcoes_agentes = {agente['nome']: agente for agente in agentes_monitoramento}
        agente_selecionado_nome = st.selectbox("Agente de monitoramento:", list(opcoes_agentes.keys()))
        agente_monitoramento = opcoes_agentes[agente_selecionado_nome]
    else:
        st.error("❌ Nenhum agente de monitoramento encontrado")
        agente_monitoramento = None
    
    if "messages_monitoramento" not in st.session_state:
        st.session_state.messages_monitoramento = []
    
    for message in st.session_state.messages_monitoramento:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
    
    if prompt := st.chat_input("Digite sua mensagem...", key="chat_monitoramento"):
        st.session_state.messages_monitoramento.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        
        with st.chat_message("assistant"):
            with st.spinner("🌱 Gerando resposta..."):
                resposta = gerar_resposta_agente(
                    prompt, st.session_state.messages_monitoramento,
                    agente_monitoramento, modelo_monitoramento, contexto_adicional
                )
                st.markdown(resposta)
                st.session_state.messages_monitoramento.append({"role": "assistant", "content": resposta})

# ==================== ABA: OTIMIZAÇÃO DE CONTEÚDO ====================
with tab_mapping["🚀 Otimização de Conteúdo"]:
    st.header("🚀 Otimização de Conteúdo")
    
    if 'conteudo_otimizado' not in st.session_state:
        st.session_state.conteudo_otimizado = None
    
    texto_para_otimizar = st.text_area("Cole o conteúdo para otimização:", height=300)
    
    col_config1, col_config2 = st.columns([2, 1])
    
    with col_config1:
        tipo_otimizacao = st.selectbox("Tipo de Otimização:", ["SEO", "Engajamento", "Conversão", "Clareza"])
    with col_config2:
        tom_voz = st.text_input("Tom de Voz:", value="Técnico")
        nivel_heading = st.selectbox("Nível de Heading:", ["H1", "H2", "H3", "H4"])
    
    usar_busca_web = st.checkbox("Usar busca web para enriquecer conteúdo", value=True)
    incluir_links_internos = st.checkbox("Incluir links internos", value=True)
    instrucoes_briefing = st.text_area("Instruções do briefing:", height=80)
    
    if st.button("🚀 Otimizar Conteúdo", type="primary", use_container_width=True):
        if texto_para_otimizar:
            with st.spinner("Processando otimização..."):
                try:
                    fontes_encontradas = ""
                    if usar_busca_web:
                        fontes_encontradas = realizar_busca_web_com_fontes(texto_para_otimizar[:500], "")
                    
                    contexto_agente = ""
                    if st.session_state.agente_selecionado:
                        contexto_agente = construir_contexto(agente_selecionado, st.session_state.segmentos_selecionados)
                    
                    prompt = f"""
                    {contexto_agente}
                    
                    ## TEXTO ORIGINAL:
                    {texto_para_otimizar}
                    
                    ## FONTES:
                    {fontes_encontradas if fontes_encontradas else "Nenhuma"}
                    
                    ## INSTRUÇÕES:
                    Tipo: {tipo_otimizacao}
                    Tom: {tom_voz}
                    Heading level: {nivel_heading}
                    Briefing: {instrucoes_briefing}
                    
                    ## REQUISITOS:
                    1. Gere 3 opções de meta title e description
                    2. Use bullets para listas quando aplicável
                    3. Todos os headings devem ser {nivel_heading}
                    4. Remova introduções genéricas
                    5. Quebre parágrafos longos
                    
                    Gere o conteúdo otimizado.
                    """
                    
                    resposta = modelo_texto.generate_content(prompt)
                    st.session_state.conteudo_otimizado = resposta.text
                    
                    st.subheader("📝 Conteúdo Otimizado")
                    st.markdown(resposta.text)
                    
                    st.download_button(
                        "💾 Baixar Conteúdo Otimizado",
                        data=resposta.text,
                        file_name=f"otimizado_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                        mime="text/plain"
                    )
                    
                except Exception as e:
                    st.error(f"Erro: {str(e)}")
        else:
            st.warning("Cole um conteúdo para otimizar")

# ==================== ABA: CRIADORA DE CALENDÁRIO ====================
with tab_mapping["📅 Criadora de Calendário"]:
    st.header("📅 Criadora de Calendário")
    
    agente = st.session_state.agente_selecionado
    st.success(f"Agente: {agente['nome']}")
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        mes_ano = st.text_input("Mês/Ano:", "FEVEREIRO 2026")
        data_inicio = st.date_input("Data início:", value=datetime.date(2026, 2, 1))
        data_fim = st.date_input("Data fim:", value=datetime.date(2026, 2, 28))
        delta_dias = (data_fim - data_inicio).days + 1
    
    with col2:
        dias_com_1_pauta = st.number_input("Dias com 1 pauta:", 0, delta_dias, 5)
        dias_com_2_pautas = st.number_input("Dias com 2 pautas:", 0, delta_dias, 15)
        dias_com_3_pautas = st.number_input("Dias com 3 pautas:", 0, delta_dias, 3)
    
    st.subheader("Produtos e Direcionais")
    produtos_direcionais = st.text_area("Produtos:", height=150)
    
    contexto_mensal = st.text_area("Contexto do mês:", height=120)
    
    if st.button("Gerar Calendário", type="primary"):
        if data_inicio >= data_fim:
            st.error("Data início deve ser anterior")
        else:
            with st.spinner("Gerando calendário..."):
                try:
                    contexto_agente = construir_contexto(agente, st.session_state.segmentos_selecionados)
                    
                    prompt = f"""
                    {contexto_agente}
                    
                    GERAR CALENDÁRIO:
                    Período: {data_inicio} a {data_fim}
                    Mês: {mes_ano}
                    
                    Dias com 1 pauta: {dias_com_1_pauta}
                    Dias com 2 pautas: {dias_com_2_pautas}
                    Dias com 3 pautas: {dias_com_3_pautas}
                    
                    Produtos: {produtos_direcionais}
                    Contexto: {contexto_mensal}
                    
                    FORMATO: CSV pronto para Excel.
                    """
                    
                    resposta = modelo_texto.generate_content(prompt)
                    calendario_csv = resposta.text
                    
                    calendario_limpo = calendario_csv.strip()
                    if '```csv' in calendario_limpo:
                        calendario_limpo = calendario_limpo.replace('```csv', '').replace('```', '')
                    
                    st.session_state.calendario_gerado = calendario_limpo
                    
                    st.subheader(f"Calendário - {mes_ano}")
                    st.text_area("CSV:", calendario_limpo, height=400)
                    
                    st.download_button(
                        "Baixar CSV",
                        data=calendario_limpo,
                        file_name=f"calendario_{mes_ano.replace(' ', '_').lower()}.csv",
                        mime="text/csv"
                    )
                    
                except Exception as e:
                    st.error(f"Erro: {str(e)}")

# ==================== ABA: PLANEJAMENTO ESTRATÉGICO ====================
with tab_mapping["📊 Planejamento Estratégico"]:
    st.header("📊 Planejamento Estratégico")
    
    col1, col2 = st.columns(2)
    
    with col1:
        nome_cliente = st.text_input('Nome do Cliente:', key="nome_cliente_planejamento")
        site_cliente = st.text_input('Site do Cliente:', key="site_cliente_planejamento")
        ramo_atuacao = st.text_input('Ramo de Atuação:', key="ramo_atuacao_planejamento")
    
    with col2:
        intuito_plano = st.text_input('Intuito do Planejamento:', key="intuito_plano_planejamento")
        publico_alvo = st.text_input('Público alvo:', key="publico_alvo_planejamento")
    
    objetivos_opcoes = [
        'Aumentar reconhecimento de marca',
        'Gerar leads qualificados',
        'Aumentar vendas',
        'Fidelizar clientes'
    ]
    objetivos_de_marca = st.selectbox('Objetivos da marca:', objetivos_opcoes, key="objetivos_marca_planejamento")
    
    referencia_da_marca = st.text_area('Referência de marca:', height=100, key="referencia_da_marca_planejamento")
    sucesso = st.text_input('O que é sucesso para a marca?:', key="sucesso_planejamento")
    concorrentes = st.text_input('Concorrentes:', key="concorrentes_planejamento")
    
    if st.button("🚀 Iniciar Planejamento Estratégico", type="primary", use_container_width=True):
        if not nome_cliente:
            st.error("Nome do cliente é obrigatório")
        else:
            with st.spinner("Gerando planejamento..."):
                try:
                    prompt = f"""
                    ## PLANEJAMENTO ESTRATÉGICO - {nome_cliente}
                    
                    Cliente: {nome_cliente}
                    Ramo: {ramo_atuacao}
                    Objetivo: {intuito_plano}
                    Público: {publico_alvo}
                    Objetivos de marca: {objetivos_de_marca}
                    Referência: {referencia_da_marca}
                    Sucesso: {sucesso}
                    Concorrentes: {concorrentes}
                    
                    Gere um planejamento estratégico completo com:
                    1. Análise SWOT
                    2. Análise PEST
                    3. Posicionamento de marca
                    4. Brand Persona
                    5. Buyer Persona
                    6. Tom de Voz
                    """
                    
                    resposta = modelo_texto.generate_content(prompt)
                    
                    st.subheader("📋 Planejamento Gerado")
                    st.markdown(resposta.text)
                    
                    st.download_button(
                        "💾 Baixar Planejamento",
                        data=resposta.text,
                        file_name=f"planejamento_{nome_cliente}_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                        mime="text/plain"
                    )
                    
                except Exception as e:
                    st.error(f"Erro: {str(e)}")

# ==================== ABA: PLANEJAMENTO DE MÍDIAS ====================
with tab_mapping["📱 Planejamento de Mídias"]:
    st.header("📱 Planejamento de Mídias e Redes")
    
    col1, col2 = st.columns(2)
    
    with col1:
        nome_cliente = st.text_input('Nome do Cliente:', key="nome_cliente_midias")
        ramo_atuacao = st.text_input('Ramo de Atuação:', key="ramo_atuacao_midias")
    
    with col2:
        intuito_plano = st.text_input('Intuito do Planejamento:', key="intuito_plano_midias")
        publico_alvo = st.text_input('Público alvo:', key="publico_alvo_midias")
    
    orcamento_total = st.number_input('Orçamento total (R$):', min_value=1000, max_value=1000000, value=10000, key="orcamento_total")
    periodo_campanha = st.selectbox('Período da campanha:', ['1 mês', '3 meses', '6 meses'], key="periodo_campanha")
    
    referencia_da_marca = st.text_area('Referência de marca:', height=100, key="referencia_da_marca_midias")
    
    if st.button("🚀 Gerar Planejamento de Mídias", type="primary", use_container_width=True):
        if not nome_cliente:
            st.error("Nome do cliente é obrigatório")
        else:
            with st.spinner("Gerando planejamento..."):
                try:
                    prompt = f"""
                    ## PLANEJAMENTO DE MÍDIAS - {nome_cliente}
                    
                    Cliente: {nome_cliente}
                    Ramo: {ramo_atuacao}
                    Objetivo: {intuito_plano}
                    Público: {publico_alvo}
                    Orçamento: R${orcamento_total}
                    Período: {periodo_campanha}
                    Referência: {referencia_da_marca}
                    
                    Gere um planejamento de mídias completo com:
                    1. Distribuição orçamentária por plataforma
                    2. Estratégia de conteúdo por pilar
                    3. Plano para Meta Ads
                    4. Plano para Google Ads
                    5. Plano para canais alternativos (TikTok, Kwai, Pinterest)
                    6. KPIs e métricas
                    """
                    
                    resposta = modelo_texto.generate_content(prompt)
                    
                    st.subheader("📋 Planejamento de Mídias")
                    st.markdown(resposta.text)
                    
                    st.download_button(
                        "💾 Baixar Planejamento",
                        data=resposta.text,
                        file_name=f"planejamento_midias_{nome_cliente}_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                        mime="text/plain"
                    )
                    
                except Exception as e:
                    st.error(f"Erro: {str(e)}")

# ==================== ABA: BRIEFING (se existir) ====================
if "📋 Briefing" in tab_mapping:
    with tab_mapping["📋 Briefing"]:
        st.header("📋 Gerador de Briefings - SYN")
        
        tab1, tab2 = st.tabs(["Briefing Individual", "Processamento em Lote (CSV)"])
        
        with tab1:
            content_input = st.text_area("Conteúdo da célula:", height=100, key="individual_content")
            data_input = st.date_input("Data prevista:", value=datetime.datetime.now(), key="individual_date")
            formato_principal = st.selectbox("Formato principal:", ["Reels + capa", "Carrossel + stories", "Blog + redes"], key="individual_format")
            
            if st.button("Gerar Briefing Individual", type="primary"):
                if content_input:
                    with st.spinner("Gerando briefing..."):
                        product, culture, action = extract_product_info(content_input)
                        
                        if product and product in PRODUCT_DESCRIPTIONS:
                            briefing = generate_briefing(content_input, product, culture, action, data_input, formato_principal)
                            st.markdown("## Briefing Gerado")
                            st.text(briefing)
                            
                            st.download_button(
                                "Baixar Briefing",
                                data=briefing,
                                file_name=f"briefing_{product}_{data_input.strftime('%Y%m%d')}.txt",
                                mime="text/plain"
                            )
                        elif product:
                            st.warning(f"Produto '{product}' não encontrado")
                        else:
                            st.error("Não foi possível identificar um produto")
        
        with tab2:
            uploaded_file = st.file_uploader("Escolha o arquivo CSV", type=['csv'])
            
            if uploaded_file:
                df = pd.read_csv(uploaded_file)
                st.success(f"CSV carregado! {len(df)} linhas")
                
                coluna_conteudo = st.selectbox("Coluna com conteúdo:", df.columns.tolist())
                data_padrao = st.date_input("Data padrão:", value=datetime.datetime.now(), key="batch_date")
                formato_padrao = st.selectbox("Formato padrão:", ["Reels + capa", "Carrossel + stories", "Blog + redes"], key="batch_format")
                
                if st.button("Processar CSV", type="primary"):
                    briefings_gerados = []
                    
                    for index, row in df.iterrows():
                        content = str(row[coluna_conteudo]) if pd.notna(row[coluna_conteudo]) else ""
                        if content:
                            product, culture, action = extract_product_info(content)
                            if product and product in PRODUCT_DESCRIPTIONS:
                                briefing = generate_briefing(content, product, culture, action, data_padrao, formato_padrao)
                                briefings_gerados.append({'produto': product, 'briefing': briefing})
                    
                    if briefings_gerados:
                        st.success(f"{len(briefings_gerados)} briefings gerados")
                        
                        for b in briefings_gerados:
                            with st.expander(f"Briefing - {b['produto']}"):
                                st.text(b['briefing'])
                    else:
                        st.warning("Nenhum briefing gerado")

# ==================== SIDEBAR - INFORMAÇÕES FINAIS ====================
with st.sidebar:
    st.markdown("---")
    st.subheader("🔐 Sistema de Isolamento")
    
    current_user = get_current_user()
    current_squad = get_current_squad()
    
    if current_squad == "admin":
        st.success("👑 **Modo Administrador**")
    else:
        st.success(f"👤 **Usuário:** {current_user.get('nome', 'Usuário')}")
        st.info(f"🏢 **Squad:** {current_squad}")
    
    agentes_usuario = listar_agentes()
    if agentes_usuario:
        categorias_count = {}
        for agente in agentes_usuario:
            cat = agente.get('categoria', 'Social')
            categorias_count[cat] = categorias_count.get(cat, 0) + 1
        
        st.markdown("### 📊 Seus Agentes")
        for categoria, count in categorias_count.items():
            st.write(f"- **{categoria}:** {count} agente(s)")
        
        st.write(f"**Total:** {len(agentes_usuario)} agente(s)")
